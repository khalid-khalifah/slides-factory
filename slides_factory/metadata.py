"""Store and read slide metadata in speaker notes.

Each rendered slide stores its template id and original JSON in the notes
text frame so doc get / slide edit can round-trip content without guessing.

Functions:
    write_metadata    — Serialize template id + data JSON into speaker notes.
    read_metadata     — Parse metadata back from speaker notes (or None).
    fallback_extract  — Best-effort title/body read when metadata is missing.
"""

import json
import logging
from typing import Any

from pptx.slide import Slide

logger = logging.getLogger(__name__)

METADATA_VERSION = 1
METADATA_PREFIX = "_sf"


def write_metadata(
    slide: Slide,
    template_id: str,
    data: dict[str, Any],
    *,
    frame_id: str | None = None,
) -> None:
    """Write template id, optional frame id, and slide JSON into speaker notes."""
    meta: dict[str, Any] = {
        "template_id": template_id,
        "version": METADATA_VERSION,
        "data": data,
    }
    if frame_id:
        meta["frame_id"] = frame_id
    payload = {METADATA_PREFIX: meta}
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = json.dumps(payload, ensure_ascii=False)


def read_metadata(slide: Slide) -> dict[str, Any] | None:
    """Read template metadata from speaker notes, or None if absent/invalid."""
    if not slide.has_notes_slide:
        return None
    text = slide.notes_slide.notes_text_frame.text.strip()
    if not text:
        return None
    try:
        payload = json.loads(text)
    except json.JSONDecodeError as exc:
        logger.warning("Corrupted slide metadata on slide %s: %s", slide.slide_id, exc)
        return None
    meta = payload.get(METADATA_PREFIX)
    if not isinstance(meta, dict):
        return None
    return meta


def fallback_extract(slide: Slide) -> dict[str, Any]:
    """Best-effort extraction when metadata is missing."""
    title = slide.shapes.title.text if slide.shapes.title else ""
    body_parts: list[str] = []
    for shape in slide.shapes:
        if not shape.has_text_frame or shape == slide.shapes.title:
            continue
        for paragraph in shape.text_frame.paragraphs:
            if paragraph.text.strip():
                body_parts.append(paragraph.text.strip())
    return {"title": title, "body": body_parts}
