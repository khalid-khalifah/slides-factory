"""Text converter — render a TextBlock into any python-pptx TextFrame.

This is the low-level primitive.  It does not create shapes or grid cells;
it only writes rich-text content into an **already-existing** ``TextFrame``.

Usage::

    from slides_factory.converters.text import render_text_block

    # Into a shape's text frame
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, ...)
    render_text_block(my_block, shape, ctx, base_size_pt=14)

    # Into a textbox
    tb = slide.shapes.add_textbox(left, top, width, height)
    render_text_block(my_block, tb, ctx, base_size_pt=14)

    # Into a table cell
    render_text_block(my_block, table.cell(0, 0), ctx, base_size_pt=12)

``text_frame`` may be a python-pptx **shape** (with ``.text_frame``) or a
bare ``TextFrame`` object — the function normalises either.
"""

from __future__ import annotations

from contextlib import suppress

from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Pt

from slides_factory.color_utils import hex_to_rgb
from slides_factory.elements.text.model import TextBlock
from slides_factory.elements.text.render import (
    _RenderParagraph,
    _RenderRun,
    prepare,
)
from slides_factory.layout.fonts import apply_shape_font
from slides_factory.render_context import RenderContext
from slides_factory.styling import theme

_ALIGN_MAP = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
    "justify": PP_ALIGN.JUSTIFY,
}


def render_text_block(
    block: TextBlock,
    text_frame: object,
    ctx: RenderContext,
    *,
    base_size_pt: float,
    base_color: str = "primary",
    base_bold: bool = False,
    alignment: str = "left",
    font_slot: str | None = None,
    vertical_anchor: str | None = None,
) -> None:
    """Render a ``TextBlock`` into an existing python-pptx ``TextFrame``.

    Params
    ------
    block:
        The rich-text document tree to render.
    text_frame:
        Any python-pptx object with a ``.text_frame`` attribute
        (shape, textbox, table cell, connector, placeholder, …), or
        a bare ``TextFrame`` instance.
    ctx:
        Render context for theme and colour-token resolution.
    base_size_pt:
        Default font size in points for paragraphs that don't set one.
    base_color:
        Theme colour token or hex string for the default text colour.
    base_bold:
        Default bold state.
    alignment:
        Default paragraph alignment (``"left"``, ``"center"``, ``"right"``,
        ``"justify"``).
    font_slot:
        Optional brand font slot name (``"body"``, ``"title"``, …) to apply
        via ``apply_shape_font``.
    vertical_anchor:
        Optional vertical anchor (``"top"``, ``"middle"``, ``"bottom"``).
    """
    # Normalise: accept either a shape (with .text_frame) or a bare TextFrame.
    if hasattr(text_frame, "text_frame"):
        shape_obj = text_frame
        tf = text_frame.text_frame
    else:
        shape_obj = None
        tf = text_frame

    tf.word_wrap = True

    if vertical_anchor is not None:
        anchor_map = {
            "top": MSO_ANCHOR.TOP,
            "middle": MSO_ANCHOR.MIDDLE,
            "bottom": MSO_ANCHOR.BOTTOM,
        }
        tf.vertical_anchor = anchor_map.get(vertical_anchor, MSO_ANCHOR.TOP)

    # Resolve the rich-text tree into render-ready paragraphs.
    render_paragraphs = prepare(
        block,
        ctx,
        base_size_pt=base_size_pt,
        base_color=base_color,
        base_bold=base_bold,
        alignment=alignment,
    )

    if not render_paragraphs:
        render_paragraphs = [_RenderParagraph(runs=[_RenderRun(text="")])]

    for index, rp in enumerate(render_paragraphs):
        para = tf.paragraphs[0] if index == 0 else tf.add_paragraph()

        if rp.alignment and rp.alignment in _ALIGN_MAP:
            para.alignment = _ALIGN_MAP[rp.alignment]

        if rp.indent_level > 0:
            para.level = rp.indent_level

        for run_idx, rr in enumerate(rp.runs):
            run = para.runs[0] if (run_idx == 0 and para.runs) else para.add_run()
            run.text = rr.text

            if rr.bold is not None:
                run.font.bold = rr.bold
            if rr.italic is not None:
                run.font.italic = rr.italic
            if rr.color_hex is not None:
                run.font.color.rgb = hex_to_rgb(rr.color_hex)
            if rr.size_pt is not None:
                run.font.size = Pt(rr.size_pt)
            if rr.strikethrough is not None:
                run.font.strikethrough = rr.strikethrough
            if rr.underline is not None:
                run.font.underline = rr.underline
            if rr.link:
                run.hyperlink.address = rr.link

    # Apply brand font if a slot was specified.
    if font_slot is not None:
        target = shape_obj if shape_obj is not None else tf
        with suppress(AttributeError):
            apply_shape_font(target, ctx, font_slot)
