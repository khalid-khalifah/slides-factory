"""SVG converter — render SVG as native python-pptx shapes.

Wraps ``svg2pptx`` to render SVG content as native, editable PowerPoint
shapes on any slide.  Unlike embedding raster images, the result is
fully editable vector artwork.

Usage::

    from slides_factory.converters.svg import render_svg_string, render_svg_file

    # Render an SVG string onto a slide at a given box
    render_svg_string('''<svg ...>...</svg>''', slide, box)

    # Render from a file
    render_svg_file("icon.svg", slide, box)
"""

from __future__ import annotations

from pathlib import Path
from typing import Union

from pptx.slide import Slide

from slides_factory.geometry import Box


def _svg_intrinsic_size(
    svg_content: str,
) -> tuple[float, float] | None:
    """Quick-parse SVG to get intrinsic width/height without full parsing.

    Returns ``(width_px, height_px)`` or ``None`` if neither the SVG
    element nor a viewBox defines dimensions.
    """
    import xml.etree.ElementTree as ET

    # Register SVG-related namespace prefixes so ElementTree doesn't
    # choke on ``xmlns:xlink`` and similar bound prefixes.
    ET.register_namespace("xlink", "http://www.w3.org/1999/xlink")
    ET.register_namespace("xml", "http://www.w3.org/XML/1998/namespace")

    root = ET.fromstring(svg_content)
    tag = root.tag if "}" not in root.tag else root.tag.split("}", 1)[1]
    if tag != "svg":
        return None

    # Try viewBox first, then width/height attributes.
    vb = root.get("viewBox")
    if vb:
        parts = vb.strip().split()
        if len(parts) == 4:
            w, h = float(parts[2]), float(parts[3])
            if w > 0 and h > 0:
                return w, h

    raw_w = root.get("width")
    raw_h = root.get("height")
    if raw_w and raw_h:
        try:
            return float(raw_w), float(raw_h)
        except ValueError:
            pass

    return None


def _fit_scale(
    svg_w: float,
    svg_h: float,
    box: Box,
) -> tuple[float, int, int]:
    """Compute scale, x-offset, and y-offset to fit SVG inside *box*
    maintaining aspect ratio (``contain`` fit, centred)."""
    from svg2pptx.geometry.units import px_to_emu

    svg_w_emu = px_to_emu(svg_w)
    svg_h_emu = px_to_emu(svg_h)

    if svg_w_emu <= 0 or svg_h_emu <= 0:
        return 1.0, box.left, box.top

    scale_x = box.width / svg_w_emu
    scale_y = box.height / svg_h_emu
    scale = min(scale_x, scale_y)

    # Centre the result in the box.
    final_w = int(svg_w_emu * scale)
    final_h = int(svg_h_emu * scale)
    offset_x = box.left + (box.width - final_w) // 2
    offset_y = box.top + (box.height - final_h) // 2

    return scale, offset_x, offset_y


def _sanitise_svg_paths(root) -> None:
    """Expand compact SVG path ``d`` attributes into clean, absolute paths.

    ``svgpathtools.parse_path().d()`` handles relative commands, shorthand
    notation, implicit repeats, and compressed coordinates (e.g. ``10-20``).
    The output is fully explicit with proper spacing — safe for svg2pptx's
    regex-based path parser.
    """
    from svgpathtools import parse_path

    SVG_NS_ = "http://www.w3.org/2000/svg"
    for path_el in root.iter(f"{{{SVG_NS_}}}path"):
        d = path_el.get("d", "").strip()
        if not d:
            continue
        try:
            parsed = parse_path(d)
            clean = parsed.d()
            if clean and clean != d:
                path_el.set("d", clean)
        except Exception:
            pass  # keep original d on parse failure


def _inline_svg_css(svg_content: str) -> str:
    """Inline CSS class styles into SVG element attributes.

    ``svg2pptx`` does not resolve CSS class definitions from ``<style>``
    blocks.  This function extracts CSS rules, resolves them onto each
    element's ``class`` attribute as inline ``fill``/``stroke``/etc.
    attributes, and removes the ``<style>`` elements.

    Existing inline attributes on elements take priority over CSS classes.
    Uses ``lxml`` (a project dependency via python-pptx) for namespace-safe
    SVG parsing and serialisation.
    """
    import re
    from lxml import etree

    SVG_NS = "http://www.w3.org/2000/svg"

    # 1. Extract CSS rules from <style> blocks.
    css_rules: dict[str, dict[str, str]] = {}
    for style_match in re.finditer(
        r"""<style[^>]*>(.*?)</style>""", svg_content, re.DOTALL
    ):
        css_text = style_match.group(1)
        for rule_match in re.finditer(
            r"""\.([-\w]+)\s*\{([^}]+)\}""", css_text
        ):
            class_name = rule_match.group(1)
            declarations = rule_match.group(2)
            styles: dict[str, str] = {}
            for prop in declarations.split(";"):
                prop = prop.strip()
                if ":" in prop:
                    key, val = prop.split(":", 1)
                    styles[key.strip()] = val.strip()
            css_rules[class_name] = styles

    if not css_rules:
        return svg_content

    # 2. Parse SVG with lxml (handles default namespaces correctly).
    parser = etree.XMLParser(remove_blank_text=True)
    root = etree.fromstring(svg_content.encode("utf-8"), parser)

    # 3. Resolve class attributes to inline fill/stroke/opacity/etc.
    svg_props = {"fill", "stroke", "stroke-width", "stroke-linecap",
                 "stroke-linejoin", "opacity", "fill-opacity",
                 "stroke-opacity"}

    for el in root.iter():
        class_attr = el.get("class") or ""
        classes = class_attr.split()
        if not classes:
            continue

        merged: dict[str, str] = {}
        for cls_name in classes:
            if cls_name in css_rules:
                merged.update(css_rules[cls_name])

        if not merged:
            continue

        for key, value in merged.items():
            if key in svg_props and el.get(key) is None:
                el.set(key, value)

    # 4. Add matching hairline stroke to filled elements to eliminate
    #    sub-pixel gaps between adjacent shapes during EMU conversion.
    for el in root.iter():
        fill = el.get("fill") or ""
        stroke = el.get("stroke") or ""
        local_tag = el.tag.split("}", 1)[-1] if "}" in el.tag else el.tag
        if local_tag in ("svg", "defs", "style", "linearGradient",
                        "radialGradient", "clipPath", "mask"):
            continue
        if fill and fill != "none" and (not stroke or stroke == "none"):
            el.set("stroke", fill)
            el.set("stroke-width", "0.5")
            el.set("stroke-linejoin", "round")
            el.set("stroke-linecap", "round")

    # 4. Remove <style> elements (lxml namespace: {svg}style).
    # 5.  Sanitise path ``d`` attributes — expand compact / minified
    #     notation (relative commands, ``10-20`` shorthand, implicit
    #     repeats) into clean, absolute, well-spaced paths that
    #     svg2pptx's regex-based parser can handle.
    _sanitise_svg_paths(root)

    ns_svg_style = f"{{{SVG_NS}}}style"
    for el in list(root.iter()):
        if el.tag == ns_svg_style:
            parent = el.getparent()
            if parent is not None:
                parent.remove(el)

    # 6.  Serialise back to string, keeping original XML declaration.
    result_bytes = etree.tostring(
        root, encoding="unicode", xml_declaration=False
    )
    # Preserve the original <svg ...> opening tag so xmlns:* attrs are
    # kept exactly as authored (lxml may reorder them).
    svg_open_end = svg_content.index(">")
    return svg_content[: svg_open_end + 1] + "\n" + result_bytes.split(">", 1)[1]


def render_svg_string(
    svg_content: str,
    slide: Slide,
    box: Box,
    *,
    scale: float | None = None,
    curve_tolerance: float = 0.01,
) -> None:
    """Render an SVG string as native PowerPoint shapes inside *box*.

    CSS class definitions from ``<style>`` blocks are automatically
    inlined before rendering (``svg2pptx`` does not resolve CSS classes).

    When *scale* is ``None`` (default), the SVG is auto-scaled to fit
    the box with ``contain`` aspect-ratio preservation and centred inside
    the box.

    When *scale* is given, it is used directly and the SVG is placed at
    the box's top-left corner.

    *curve_tolerance* controls the smoothness of bezier curve approximation
    (in SVG viewBox units).  Lower values = smoother curves = more shapes.
    Default 0.01 produces ~35K segments for a typical logo; increase for
    very large SVGs to reduce shape count, decrease for max smoothness.
    """
    from svg2pptx import Config, SVGConverter

    # Inline CSS class styles before passing to svg2pptx.
    svg_content = _inline_svg_css(svg_content)

    size = _svg_intrinsic_size(svg_content)

    if scale is not None:
        final_scale = scale
        offset_x = box.left
        offset_y = box.top
    elif size is not None:
        svg_w, svg_h = size
        final_scale, offset_x, offset_y = _fit_scale(svg_w, svg_h, box)
    else:
        # No size info — use 1.0 and place at box origin.
        final_scale = 1.0
        offset_x = box.left
        offset_y = box.top

    config = Config(
        scale=final_scale,
        offset_x=offset_x,
        offset_y=offset_y,
        curve_tolerance=curve_tolerance,
        disable_shadows=True,
    )
    converter = SVGConverter(config=config)
    converter.add_string_to_slide(svg_content, slide)


def render_svg_file(
    path: Union[str, Path],
    slide: Slide,
    box: Box,
    *,
    scale: float | None = None,
    curve_tolerance: float = 0.01,
    method: str = "auto",
) -> None:
    """Render an SVG file onto a slide.

    *method* controls the rendering engine:

    ============= ========================================================
    ``"auto"``   EMF (Inkscape) → PNG sibling → svg2pptx (default)
    ``"emf"``    Force Inkscape EMF conversion (native editable vectors)
    ``"png"``    Force PNG raster (``slide.shapes.add_picture``)
    ``"svg2pptx"`` Force svg2pptx freeform conversion
    ============= ========================================================

    Inkscape must be installed for ``"emf"`` (``brew install --cask inkscape``
    on macOS).  EMF shapes can be ungrouped in PowerPoint to get native
    editable geometry with correct fills.
    """
    svg_path = Path(path)

    if method in ("auto", "emf") and _inkscape_available():
        _render_via_inkscape(svg_path, slide, box)
        return

    if method in ("auto", "png"):
        png_path = svg_path.with_suffix(".png")
        if png_path.is_file():
            slide.shapes.add_picture(
                str(png_path), box.left, box.top, box.width, box.height
            )
            return
        if method == "png":
            raise FileNotFoundError(
                f"PNG sibling not found for {svg_path} (method='png')"
            )

    # Fallback: svg2pptx
    content = svg_path.read_text(encoding="utf-8")
    render_svg_string(content, slide, box, scale=scale, curve_tolerance=curve_tolerance)


def _inkscape_available() -> bool:
    """True when Inkscape CLI is installed and reachable."""
    import shutil
    return shutil.which("inkscape") is not None


def _render_via_inkscape(svg_path: Path, slide: Slide, box: Box) -> None:
    """Convert SVG to EMF via Inkscape and embed as an editable picture."""
    import subprocess
    import tempfile

    emf_path = Path(tempfile.mktemp(suffix=".emf"))
    try:
        result = subprocess.run(
            [
                "inkscape",
                str(svg_path),
                f"--export-filename={emf_path}",
            ],
            capture_output=True,
            text=True,
            timeout=30,
        )
        if result.returncode != 0 or not emf_path.is_file():
            raise RuntimeError(
                f"Inkscape EMF conversion failed: {result.stderr.strip()}"
            )
        slide.shapes.add_picture(
            str(emf_path), box.left, box.top, box.width, box.height
        )
    finally:
        if emf_path.exists():
            emf_path.unlink()
