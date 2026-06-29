"""Rich-text document model, DSL, HTML parser, render pipeline, and converter.

Usage::

    from slides_factory.converters.text import TextBlock, text, parse_html, render_text_block

    tb = slide.shapes.add_textbox(left, top, width, height)
    render_text_block(parse_html("<b>Hello</b>"), tb, ctx, base_size_pt=14)
"""

from contextlib import suppress

from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Pt

from slides_factory.color_utils import hex_to_rgb
from slides_factory.converters.text.dsl import text
from slides_factory.converters.text.html import parse_html
from slides_factory.converters.text.model import (
    ListItem,
    ListStyle,
    Paragraph,
    TextBlock,
    TextRun,
)
from slides_factory.converters.text.render import (
    _RenderParagraph,
    _RenderRun,
    prepare,
)
from slides_factory.layout.fonts import apply_shape_font
from slides_factory.render_context import RenderContext
from slides_factory.styling import theme

_ALIGN_MAP = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
    "justify": PP_ALIGN.JUSTIFY,
}


def render_text_block(
    block: TextBlock,
    text_frame: object,
    ctx: RenderContext,
    *,
    base_size_pt: float | None = None,
    base_color: str | None = None,
    base_bold: bool | None = None,
    alignment: str | None = None,
    font_slot: str | None = None,
    vertical_anchor: str | None = None,
) -> None:
    """Render a ``TextBlock`` into an existing python-pptx ``TextFrame``.

    ``text_frame`` may be a python-pptx **shape** (with ``.text_frame``) or
    a bare ``TextFrame`` object — the function normalises either.

    Params
    ------
    block:
        The rich-text document tree to render.
    text_frame:
        Any python-pptx object with a ``.text_frame`` attribute
        (shape, textbox, table cell, connector, placeholder, …), or
        a bare ``TextFrame`` instance.
    ctx:
        Render context for theme and colour-token resolution.
    base_size_pt:
        Default font size in points for paragraphs that don't set one.
    base_color:
        Theme colour token or hex string for the default text colour.
    base_bold:
        Default bold state.
    alignment:
        Default paragraph alignment.
    font_slot:
        Optional brand font slot to apply.
    vertical_anchor:
        Optional vertical anchor (``"top"``, ``"middle"``, ``"bottom"``).
    """
    if hasattr(text_frame, "text_frame"):
        shape_obj = text_frame
        tf = text_frame.text_frame
    else:
        shape_obj = None
        tf = text_frame

    tf.word_wrap = True

    if vertical_anchor is not None:
        anchor_map = {
            "top": MSO_ANCHOR.TOP,
            "middle": MSO_ANCHOR.MIDDLE,
            "bottom": MSO_ANCHOR.BOTTOM,
        }
        tf.vertical_anchor = anchor_map.get(vertical_anchor, MSO_ANCHOR.TOP)

    # ── Resolve defaults from the block itself ───────────────────────────
    # Block-level fields (set by <div>) act as fallback when the caller
    # doesn't explicitly provide the corresponding function parameter.

    if base_size_pt is None and block.font_size is not None:
        base_size_pt = block.font_size
    if base_color is None and block.color is not None:
        base_color = block.color
    if base_bold is None and block.bold is not None:
        base_bold = block.bold
    if alignment is None and block.align is not None:
        alignment = block.align
    if font_slot is None and block.font_family is not None:
        font_slot = block.font_family

    # Hardcoded fallbacks if everything is None.
    if base_size_pt is None:
        base_size_pt = 14.0
    if base_color is None:
        base_color = "primary"
    if base_bold is None:
        base_bold = False
    if alignment is None:
        alignment = "left"

    render_paragraphs = prepare(
        block,
        ctx,
        base_size_pt=base_size_pt,
        base_color=base_color,
        base_bold=base_bold,
        alignment=alignment,
    )

    if not render_paragraphs:
        render_paragraphs = [_RenderParagraph(runs=[_RenderRun(text="")])]

    for index, rp in enumerate(render_paragraphs):
        para = tf.paragraphs[0] if index == 0 else tf.add_paragraph()

        if rp.alignment and rp.alignment in _ALIGN_MAP:
            para.alignment = _ALIGN_MAP[rp.alignment]

        if rp.indent_level > 0:
            para.level = rp.indent_level

        # Native PPTX list bullet — set via XML on the paragraph.
        if rp.bullet_type is not None:
            from lxml import etree
            from pptx.oxml.ns import qn

            pPr = para._pPr  # noqa: SLF001
            if pPr is None:
                pPr = para._p  # noqa: SLF001
                # create pPr element
                pPr = etree.SubElement(pPr, qn("a:pPr"))
            # Remove any existing buChar / buNone elements.
            for child in list(pPr):
                if child.tag in (qn("a:buChar"), qn("a:buNone")):
                    pPr.remove(child)
            # Add the bullet character.
            bullet_chars = {
                "disc": "\u2022",
                "circle": "\u25E6",
                "square": "\u25AA",
                "decimal": "1.",
                "hyphen": "\u2013",
                "none": None,
            }
            char = bullet_chars.get(rp.bullet_type)
            if char is not None:
                bu = etree.SubElement(pPr, qn("a:buChar"))
                bu.set("char", char)
            else:
                etree.SubElement(pPr, qn("a:buNone"))

        for run_idx, rr in enumerate(rp.runs):
            run = para.runs[0] if (run_idx == 0 and para.runs) else para.add_run()
            run.text = rr.text

            if rr.bold is not None:
                run.font.bold = rr.bold
            if rr.italic is not None:
                run.font.italic = rr.italic
            if rr.color_hex is not None:
                run.font.color.rgb = hex_to_rgb(rr.color_hex)
            if rr.size_pt is not None:
                run.font.size = Pt(rr.size_pt)
            if rr.strikethrough is not None:
                run.font.strikethrough = rr.strikethrough
            if rr.underline is not None:
                run.font.underline = rr.underline
            if rr.link:
                run.hyperlink.address = rr.link

    if font_slot is not None:
        target = shape_obj if shape_obj is not None else tf
        with suppress(AttributeError):
            apply_shape_font(target, ctx, font_slot)


__all__ = [
    "ListItem",
    "ListStyle",
    "Paragraph",
    "TextBlock",
    "TextRun",
    "_RenderParagraph",
    "_RenderRun",
    "parse_html",
    "prepare",
    "render_text_block",
    "text",
]
