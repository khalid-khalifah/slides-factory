"""RTL layout and Arabic text styling helpers used inside templates.

Classes:
    RTLLayout — Mirrors shape positions and applies RTL paragraph styling.

Methods:
    x                    — Compute mirrored left coordinate for custom shapes.
    position_for_reading — Mirror a placeholder/shape and set RTL text direction.
    style_label_box      — Apply centered RTL styling to a label inside a shape.
    mirror_row           — Mirror a horizontal row of shape left positions for RTL.
    _set_left            — Move shape horizontally without resetting top/width/height.
    _apply_paragraph_rtl — Set rtl=1, alignment, font language on one paragraph.
    _swap_margins        — Swap left/right text insets after mirroring body text.
    _style_text_frame    — Apply rtlCol and RTL styling to all paragraphs.
    _style_shape         — Apply RTL text styling to a shape's text frame.
"""

from __future__ import annotations

from typing import TYPE_CHECKING

from pptx.enum.text import PP_ALIGN
from pptx.util import Length

if TYPE_CHECKING:
    from pptx.text.text import TextFrame

    from slides_factory.render_context import RenderContext


class RTLLayout:
    """Per-template RTL positioning and text styling."""

    def __init__(self, ctx: RenderContext):
        self.ctx = ctx

    def x(self, left: int | Length, width: int | Length) -> int:
        """Compute RTL-aware left coordinate: slide_width - left - width."""
        return self.ctx.mirror_left(left, width)

    @staticmethod
    def _set_left(shape, left: int) -> None:
        """Set horizontal position without resetting other geometry.

        python-pptx clears ``top`` when ``left`` is assigned on placeholders.
        """
        top, width, height = shape.top, shape.width, shape.height
        shape.left = left
        shape.top = top
        shape.width = width
        shape.height = height

    def position_for_reading(
        self,
        shape,
        *,
        align: PP_ALIGN = PP_ALIGN.RIGHT,
        swap_margins: bool = False,
    ) -> None:
        """Mirror shape position and apply RTL text direction."""
        if not self.ctx.rtl:
            return
        self._set_left(shape, self.ctx.mirror_left(shape.left, shape.width))
        self._style_shape(shape, alignment=align, swap_margins=swap_margins)

    def _apply_paragraph_rtl(self, paragraph, *, alignment: PP_ALIGN) -> None:
        """Set rtl flag and alignment on a single paragraph."""
        p_pr = paragraph._p.get_or_add_pPr()
        p_pr.set("rtl", "1")
        if paragraph.alignment != PP_ALIGN.CENTER:
            paragraph.alignment = alignment

    def _swap_margins(self, text_frame: TextFrame) -> None:
        """Swap left/right text insets so padding follows the mirrored layout."""
        body_pr = text_frame._bodyPr
        body_pr.lIns, body_pr.rIns = body_pr.rIns, body_pr.lIns

    def _style_text_frame(
        self,
        text_frame: TextFrame,
        *,
        alignment: PP_ALIGN = PP_ALIGN.RIGHT,
        swap_margins: bool = False,
    ) -> None:
        """Apply rtlCol and RTL paragraph styling to every paragraph in a text frame."""
        if not self.ctx.rtl:
            return
        text_frame._bodyPr.set("rtlCol", "1")
        if swap_margins:
            self._swap_margins(text_frame)
        for paragraph in text_frame.paragraphs:
            self._apply_paragraph_rtl(paragraph, alignment=alignment)

    def _style_shape(
        self,
        shape,
        *,
        alignment: PP_ALIGN = PP_ALIGN.RIGHT,
        swap_margins: bool = False,
    ) -> None:
        """Apply RTL text styling to a shape that has a text frame."""
        if not shape.has_text_frame:
            return
        self._style_text_frame(
            shape.text_frame,
            alignment=alignment,
            swap_margins=swap_margins,
        )

    def style_label_box(self, shape) -> None:
        """Style label text inside a colored box (centered, RTL when active)."""
        self._style_shape(shape, alignment=PP_ALIGN.CENTER)

    def mirror_row(
        self,
        ltr_lefts: list[int],
        box_width: int | Length,
        *,
        min_margin: int = 91440,
    ) -> list[int]:
        """Mirror left positions for a horizontal row of equally-sized shapes."""
        width_i = int(box_width)
        if not self.ctx.rtl:
            return ltr_lefts

        span_left = ltr_lefts[0] if ltr_lefts else 0
        span_right = int(ltr_lefts[-1] + width_i) if ltr_lefts else span_left
        slide_width = self.ctx.slide_width
        mirrored: list[int] = []
        for ltr_left in ltr_lefts:
            left = span_right - (ltr_left - span_left) - width_i
            left = max(min_margin, min(left, slide_width - width_i - min_margin))
            mirrored.append(left)
        return mirrored
