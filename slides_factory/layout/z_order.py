"""Send shapes behind placeholders without breaking OOXML child order."""

from __future__ import annotations

from pptx.slide import Slide

# Index in p:spTree — after nvGrpSpPr (0) and grpSpPr (1).
BACKGROUND_Z_ORDER = 2


def send_shape_behind_content(slide: Slide, shape, *, z_order: int = BACKGROUND_Z_ORDER) -> None:
    """Move *shape* behind placeholders on *slide*."""
    sp_tree = slide.shapes._spTree
    element = shape._element
    sp_tree.remove(element)
    sp_tree.insert(z_order, element)
