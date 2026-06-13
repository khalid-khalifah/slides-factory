"""OOXML shape locking for frame chrome (background, logo, footer, etc.)."""

from __future__ import annotations

from lxml import etree
from pptx.shapes.base import BaseShape
from pptx.slide import Slide

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
NSMAP = {"a": A_NS, "p": P_NS}

_MOVE_LOCKS = {
    "noGrp": "1",
    "noRot": "1",
    "noMove": "1",
    "noResize": "1",
}


def lock_shape(shape: BaseShape) -> None:
    """Lock a shape so it is harder to move or resize in PowerPoint."""
    tag = shape._element.tag.split("}")[-1]
    if tag == "pic":
        _lock_picture(shape._element)
    elif tag == "sp":
        _lock_auto_shape(shape._element, text=shape.has_text_frame)


def lock_shapes_added(slide: Slide, existing: set[int]) -> None:
    """Lock shapes added since *existing* snapshot of ``id(shape._element)``."""
    for shape in slide.shapes:
        if id(shape._element) not in existing:
            lock_shape(shape)


def _apply_lock_attrs(locks: etree._Element, attrs: dict[str, str]) -> None:
    for key, value in attrs.items():
        locks.set(key, value)


def _mark_user_drawn(element: etree._Element) -> None:
    """Mark shape as user-drawn — PowerPoint honors locks more reliably."""
    tag = element.tag.split("}")[-1]
    if tag == "pic":
        nv = element.find("p:nvPicPr", NSMAP)
    elif tag == "sp":
        nv = element.find("p:nvSpPr", NSMAP)
    else:
        return
    if nv is None:
        return
    nv_pr = nv.find("p:nvPr", NSMAP)
    if nv_pr is None:
        nv_pr = etree.SubElement(nv, f"{{{P_NS}}}nvPr")
    nv_pr.set("userDrawn", "1")


def _lock_picture(element: etree._Element) -> None:
    c_nv = element.find(".//p:cNvPicPr", NSMAP)
    if c_nv is None:
        return
    c_nv.set("preferRelativeResize", "0")
    locks = c_nv.find("a:picLocks", NSMAP)
    if locks is None:
        locks = etree.SubElement(c_nv, f"{{{A_NS}}}picLocks")
    _apply_lock_attrs(locks, _MOVE_LOCKS)
    if "noChangeAspect" not in locks.attrib:
        locks.set("noChangeAspect", "1")
    _mark_user_drawn(element)


def _lock_auto_shape(element: etree._Element, *, text: bool) -> None:
    c_nv = element.find(".//p:cNvSpPr", NSMAP)
    if c_nv is None:
        return
    locks = c_nv.find("a:spLocks", NSMAP)
    if locks is None:
        locks = etree.SubElement(c_nv, f"{{{A_NS}}}spLocks")
    _apply_lock_attrs(locks, _MOVE_LOCKS)
    if text:
        locks.set("noTextEdit", "1")
    _mark_user_drawn(element)
