"""Apply brand theme fonts to slide text.

Functions:
    font_family_from_file        — Read typographic family name from a TTF/OTF file.
    language_for_locale          — Map BCP-47 locale tag to PowerPoint language id.
    resolve_font_name            — Return font family for a theme slot (falls back to Arial).
    _uses_complex_script         — True when RTL or complex-script locale is active.
    _set_complex_script_typeface — Set ``a:cs`` typeface for Arabic/RTL glyphs.
    apply_paragraph_font         — Set font family and language on every run in a paragraph.
    apply_text_frame_font        — Apply a font slot to all paragraphs in a text frame.
    apply_shape_font             — Apply a font slot to a shape's text frame when present.
"""

from __future__ import annotations

from pathlib import Path
from typing import TYPE_CHECKING, Literal

from fontTools.ttLib import TTFont
from lxml import etree
from pptx.enum.lang import MSO_LANGUAGE_ID

_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"

if TYPE_CHECKING:
    from pptx.text.text import TextFrame

    from slides_factory.render_context import RenderContext

FontSlot = Literal["title", "body", "footer"]


def font_family_from_file(path: Path | str) -> str:
    """Read the typographic family name embedded in a TTF/OTF file."""
    font = TTFont(path)
    names = {record.nameID: record.toUnicode() for record in font["name"].names}
    for name_id in (1, 4, 16):
        value = names.get(name_id)
        if value:
            return value
    raise ValueError(f"font family name not found in {path}")


def language_for_locale(locale: str) -> MSO_LANGUAGE_ID:
    """Map a BCP-47 locale tag to a PowerPoint language id."""
    if locale.startswith("ar"):
        return MSO_LANGUAGE_ID.ARABIC
    if locale.startswith("he"):
        return MSO_LANGUAGE_ID.HEBREW
    if locale.startswith("fa"):
        return MSO_LANGUAGE_ID.FARSI
    if locale.startswith("ur"):
        return MSO_LANGUAGE_ID.URDU
    return MSO_LANGUAGE_ID.ENGLISH_US


def resolve_font_name(ctx: RenderContext, slot: FontSlot = "body") -> str:
    """Return the font family for a theme slot, falling back to Arial without a brand."""
    if ctx.brand is None:
        return ctx.font_name
    return ctx.brand.fonts.family_for(ctx.brand, slot)


def _uses_complex_script(ctx: RenderContext) -> bool:
    return ctx.rtl or ctx.locale.startswith(("ar", "fa", "he", "ur"))


def _set_complex_script_typeface(r_pr, font_name: str) -> None:
    """Set ``a:cs`` typeface — PowerPoint uses this slot for Arabic/RTL glyphs."""
    cs_tag = f"{{{_A_NS}}}cs"
    latin_tag = f"{{{_A_NS}}}latin"
    cs = r_pr.find(cs_tag)
    if cs is None:
        cs = etree.Element(cs_tag)
        cs.set("typeface", font_name)
        latin = r_pr.find(latin_tag)
        if latin is not None:
            latin.addnext(cs)
        else:
            r_pr.append(cs)
    else:
        cs.set("typeface", font_name)


def apply_paragraph_font(paragraph, ctx: RenderContext, slot: FontSlot = "body") -> None:
    """Set font family and language on every run in one paragraph."""
    font_name = resolve_font_name(ctx, slot)
    language = language_for_locale(ctx.locale)
    use_cs = _uses_complex_script(ctx)
    for run in paragraph.runs:
        run.font.name = font_name
        run.font.language_id = language
        if use_cs:
            _set_complex_script_typeface(run._r.get_or_add_rPr(), font_name)


def apply_text_frame_font(
    text_frame: TextFrame,
    ctx: RenderContext,
    slot: FontSlot = "body",
) -> None:
    """Apply a single font slot to all paragraphs in a text frame."""
    for paragraph in text_frame.paragraphs:
        apply_paragraph_font(paragraph, ctx, slot)


def apply_shape_font(shape, ctx: RenderContext, slot: FontSlot = "body") -> None:
    """Apply a font slot to a shape's text frame when present."""
    if shape.has_text_frame:
        apply_text_frame_font(shape.text_frame, ctx, slot)
