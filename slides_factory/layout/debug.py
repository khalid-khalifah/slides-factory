"""Layout debug visualisation — diagnostic shapes behind rendered content.

The ``render_debug_layer()`` function draws grid lines, cell boundaries,
padding/gap regions, playground boundary, and cell-index labels so that
layout authors can verify their grid configuration visually.
"""

from __future__ import annotations

from typing import TYPE_CHECKING

from lxml import etree
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap, qn
from pptx.slide import Slide

if TYPE_CHECKING:
    from slides_factory.layout.grid import Cell, GridStyle
    from slides_factory.render_context import RenderContext


# ---------------------------------------------------------------------------
# Colour constants (RGB hex)
# ---------------------------------------------------------------------------
_PLAYGROUND_COLOR = "FF0000"  # red
_GRID_LINE_COLOR = "888888"  # gray
_PADDING_COLOR = "FFD700"  # gold
_GAP_COLOR = "DDDDDD"  # light gray
_CELL_BOUNDARY_COLOR = "3366FF"  # blue
_LABEL_COLOR = "000000"  # black


# ---------------------------------------------------------------------------
# XML helpers (python-pptx high-level API is limited for these features)
# ---------------------------------------------------------------------------


def _set_dashed_line(shape: object) -> None:
    """Set the shape's line to a dashed stroke via XML manipulation."""
    sp_pr = shape._element.spPr  # noqa: SLF001
    ln = sp_pr.find(qn("a:ln"))
    if ln is None:
        ln = etree.SubElement(sp_pr, qn("a:ln"))
        ln.set("w", "6350")  # 0.5 pt
    # Set solid fill for the line colour (already set via fore_color)
    # Add dash style
    prst_dash = etree.SubElement(ln, qn("a:prstDash"))
    prst_dash.set("val", "dash")


def _set_no_fill(shape: object) -> None:
    """Remove any fill from the shape (line-only)."""
    sp_pr = shape._element.spPr  # noqa: SLF001
    for fill_elem in sp_pr.findall(qn("a:solidFill")):
        sp_pr.remove(fill_elem)
    # Ensure noFill is set
    etree.SubElement(sp_pr, qn("a:noFill"))


def _set_fill_alpha(shape: object, alpha_pct: int) -> None:
    """Set the shape's solid fill transparency (0 = invisible, 100 = opaque)."""
    sp_pr = shape._element.spPr  # noqa: SLF001
    solid_fill = sp_pr.find(qn("a:solidFill"))
    if solid_fill is not None:
        srgb = solid_fill.find(qn("a:srgbClr"))
        if srgb is not None:
            alpha = etree.SubElement(srgb, qn("a:alpha"))
            alpha.set("val", str(int(alpha_pct * 1000)))  # in 1/1000ths of percent


def _set_shape_outline(shape: object, hex_color: str, width_emu: int = 6350) -> None:
    """Apply a solid outline (border) to a shape.

    *hex_color* is 6 hex chars without ``#``.
    """
    sp_pr = shape._element.spPr  # noqa: SLF001
    ln = sp_pr.find(qn("a:ln"))
    if ln is None:
        ln = etree.SubElement(sp_pr, qn("a:ln"))
    ln.set("w", str(width_emu))
    # Solid fill with colour
    solid_fill = ln.find(qn("a:solidFill"))
    if solid_fill is None:
        solid_fill = etree.SubElement(ln, qn("a:solidFill"))
    else:
        for old in solid_fill.iterchildren():
            solid_fill.remove(old)
    srgb = etree.SubElement(solid_fill, qn("a:srgbClr"))
    srgb.set("val", hex_color)


def _set_shape_fill_color(shape: object, hex_color: str, alpha_pct: int = 20) -> None:
    """Set solid fill + colour on a shape with optional transparency."""
    shape.fill.solid()
    from pptx.dml.color import RGBColor

    shape.fill.fore_color.rgb = RGBColor(*bytes.fromhex(hex_color))
    _set_fill_alpha(shape, alpha_pct)


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------


def render_debug_layer(
    slide: Slide,
    grid: GridStyle,
    cells: list[Cell],
    *,
    ctx: RenderContext,
) -> None:
    """Draw diagnostic shapes for *grid* and *cells* into *slide*.

    All shapes are rendered *before* normal content so they appear behind it
    (python-pptx order-of-creation z-ordering).
    """
    region = ctx.playground
    if region is None:
        return

    region_left, region_top, region_w, region_h = region

    _draw_playground_boundary(slide, region)
    _draw_padding_region(slide, grid, region)
    _draw_grid_lines(slide, grid, region)
    _draw_gap_regions(slide, grid, region)
    _draw_cell_boundaries(slide, cells)
    _draw_cell_labels(slide, cells, region)


def _draw_playground_boundary(slide: Slide, region: tuple[int, int, int, int]) -> None:
    """Dashed red rectangle around the playground region."""
    left, top, width, height = region
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    _set_no_fill(shape)
    _set_shape_outline(shape, _PLAYGROUND_COLOR, width_emu=12700)  # 1 pt
    _set_dashed_line(shape)


def _draw_padding_region(
    slide: Slide, grid: GridStyle, region: tuple[int, int, int, int]
) -> None:
    """Semi-transparent gold fill between playground edge and inner grid."""
    left, top, width, height = region
    pad_x = grid.pad_x * width
    pad_y = grid.pad_y * height

    if pad_x <= 0 and pad_y <= 0:
        return

    # We draw 4 rects: top, bottom, left, right padding strips.
    _draw_rect(slide, left, top, width, pad_y, _PADDING_COLOR, alpha_pct=15)  # top
    _draw_rect(
        slide, left, top + height - pad_y, width, pad_y, _PADDING_COLOR, alpha_pct=15
    )  # bottom
    _draw_rect(slide, left, top + pad_y, pad_x, height - 2 * pad_y, _PADDING_COLOR, alpha_pct=15)
    _draw_rect(
        slide,
        left + width - pad_x,
        top + pad_y,
        pad_x,
        height - 2 * pad_y,
        _PADDING_COLOR,
        alpha_pct=15,
    )


def _draw_grid_lines(
    slide: Slide, grid: GridStyle, region: tuple[int, int, int, int]
) -> None:
    """Dashed gray lines at column/row track boundaries."""
    # We need the inner region dimensions to calculate track positions.
    left, top, width, height = region
    pad_x = grid.pad_x * width
    pad_y = grid.pad_y * height
    inner_left = left + pad_x
    inner_top = top + pad_y
    inner_w = width - 2 * pad_x
    inner_h = height - 2 * pad_y

    if inner_w <= 0 or inner_h <= 0:
        return

    ncols = len(grid.columns)
    # Calculate column track starts.
    gap_x = grid.col_gap * inner_w
    available_x = inner_w - gap_x * (ncols - 1)
    if available_x > 0:
        sum_ratios_x = sum(grid.columns)
        col_sizes = [available_x * r / sum_ratios_x for r in grid.columns]
        cursor = 0.0
        for i in range(ncols):
            cursor += col_sizes[i]
            x = int(inner_left + cursor + gap_x * i)
            # Vertical line at column boundary
            line = slide.shapes.add_connector(
                1, x, inner_top, x, inner_top + inner_h  # 1 = straight connector
            )
            _set_no_fill(line)
            _set_shape_outline(line, _GRID_LINE_COLOR, width_emu=3175)  # 0.25 pt
            _set_dashed_line(line)

    nrows = len(grid.rows)
    gap_y = grid.row_gap * inner_h
    available_y = inner_h - gap_y * (nrows - 1)
    if available_y > 0 and nrows > 0:
        sum_ratios_y = sum(grid.rows)
        row_sizes = [available_y * r / sum_ratios_y for r in grid.rows]
        cursor = 0.0
        for i in range(nrows):
            cursor += row_sizes[i]
            y = int(inner_top + cursor + gap_y * i)
            # Horizontal line at row boundary
            line = slide.shapes.add_connector(
                1, inner_left, y, inner_left + inner_w, y
            )
            _set_no_fill(line)
            _set_shape_outline(line, _GRID_LINE_COLOR, width_emu=3175)
            _set_dashed_line(line)


def _draw_gap_regions(
    slide: Slide, grid: GridStyle, region: tuple[int, int, int, int]
) -> None:
    """Light semi-transparent bands between tracks."""
    left, top, width, height = region
    pad_x = grid.pad_x * width
    pad_y = grid.pad_y * height
    inner_left = left + pad_x
    inner_top = top + pad_y
    inner_w = width - 2 * pad_x
    inner_h = height - 2 * pad_y

    if inner_w <= 0 or inner_h <= 0:
        return

    ncols = len(grid.columns)
    gap_px = int(grid.col_gap * inner_w)
    gap_py = int(grid.row_gap * inner_h)

    # Column gaps: vertical bands between columns
    if gap_px > 0 and ncols > 1:
        available_x = inner_w - grid.col_gap * inner_w * (ncols - 1)
        if available_x > 0:
            sum_ratios_x = sum(grid.columns)
            col_sizes = [available_x * r / sum_ratios_x for r in grid.columns]
            cursor = 0.0
            for i in range(ncols - 1):
                cursor += col_sizes[i]
                gx = int(inner_left + cursor + gap_px * i)
                _draw_rect(slide, gx, inner_top, gap_px, inner_h, _GAP_COLOR, alpha_pct=12)

    # Row gaps: horizontal bands between rows
    nrows = len(grid.rows)
    if gap_py > 0 and nrows > 1:
        available_y = inner_h - grid.row_gap * inner_h * (nrows - 1)
        if available_y > 0:
            sum_ratios_y = sum(grid.rows)
            row_sizes = [available_y * r / sum_ratios_y for r in grid.rows]
            cursor = 0.0
            for i in range(nrows - 1):
                cursor += row_sizes[i]
                gy = int(inner_top + cursor + gap_py * i)
                _draw_rect(slide, inner_left, gy, inner_w, gap_py, _GAP_COLOR, alpha_pct=12)


def _draw_cell_boundaries(slide: Slide, cells: list[Cell]) -> None:
    """Thin solid blue rectangle around each cell."""
    for cell in cells:
        _draw_rect(slide, cell.left, cell.top, cell.width, cell.height, _CELL_BOUNDARY_COLOR)
        _set_shape_outline(
            slide.shapes[-1], _CELL_BOUNDARY_COLOR, width_emu=3175  # 0.25 pt
        )


def _draw_cell_labels(slide: Slide, cells: list[Cell], region: object) -> None:
    """Small text ``[row,col]`` in the top-left corner of each cell."""
    from pptx.dml.color import RGBColor
    from pptx.util import Emu, Pt

    for cell in cells:
        label = f"[{cell.row},{cell.col}]"
        tb_w = Emu(min(cell.width, 914400))  # cap at 0.6"
        tb_h = Emu(457200)  # 0.5"
        tb = slide.shapes.add_textbox(cell.left, cell.top, tb_w, tb_h)
        tf = tb.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(8)
        p.font.color.rgb = RGBColor(0x00, 0x00, 0x00)


def _draw_rect(
    slide: Slide,
    left: int,
    top: int,
    width: int,
    height: int,
    hex_color: str,
    alpha_pct: int = 0,
) -> None:
    """Add a rectangle shape with optional semi-transparent fill."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    if alpha_pct > 0:
        _set_shape_fill_color(shape, hex_color, alpha_pct)
    else:
        _set_no_fill(shape)
        _set_shape_outline(shape, hex_color, width_emu=3175)
