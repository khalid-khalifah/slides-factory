"""Presentation session management — handles I/O, themes, and document shell initialization.

Classes:
    PresentationSession — Encapsulates the python-pptx Presentation object and its lifecycle.
"""

from __future__ import annotations

from pathlib import Path
from typing import Any

from pptx import Presentation

from slides_factory.brand import BrandTheme, load_brand
from slides_factory.brand.doc import get_document_brand_path, set_document_brand
from slides_factory.locale import set_document_settings


class PresentationSession:
    """Encapsulates a PowerPoint presentation and its session-level operations."""

    def __init__(self, prs: Presentation):
        self.prs = prs

    @property
    def presentation(self) -> Presentation:
        """The underlying python-pptx Presentation object."""
        return self.prs

    @classmethod
    def open(cls, path: Path) -> PresentationSession:
        """Load an existing .pptx presentation from disk."""
        return cls(Presentation(str(path)))

    @classmethod
    def create_new(
        cls,
        theme: Path | None = None,
        brand: Path | None = None,
        rtl: bool = False,
        locale: str = "en",
        default_theme_provider: Any = None,
        delete_slide_callback: Any = None,
    ) -> PresentationSession:
        """Create an in-memory presentation shell."""
        brand_theme: BrandTheme | None = None
        if brand is not None:
            brand_theme = load_brand(brand)

        if theme is not None:
            theme_path = theme
        elif brand_theme and brand_theme.base_pptx is not None:
            theme_path = brand_theme.resolve_base_pptx()
        else:
            # Note: default_theme_provider should be a function that returns the Path to the default theme
            theme_path = (
                default_theme_provider() if default_theme_provider else Path("default.pptx")
            )

        prs = Presentation(str(theme_path))
        if brand_theme is not None:
            brand_theme.apply_page_size(prs)

        # Clear default slides
        while len(prs.slides) > 0:
            if delete_slide_callback:
                delete_slide_callback(prs, 0)
            else:
                # Fallback if no callback provided (though in this project we have a specific way to delete)
                pass

        set_document_settings(prs, rtl=rtl, locale=locale)
        if brand is not None:
            set_document_brand(prs, brand.resolve())

        return cls(prs)

    def save(self, path: Path) -> None:
        """Persist the presentation to disk."""
        path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(path))

    def load_brand(self) -> BrandTheme | None:
        """Load brand YAML stored on the document, if configured."""
        brand_path = get_document_brand_path(self.prs)
        if brand_path is None or not brand_path.is_file():
            return None
        return load_brand(brand_path)

    def update_rtl(self, rtl: bool, locale: str | None = None) -> dict[str, str | bool]:
        """Update document RTL/locale flags."""
        from slides_factory.locale import get_document_locale

        active_locale = locale or get_document_locale(self.prs)
        set_document_settings(self.prs, rtl=rtl, locale=active_locale)
        return {"rtl": rtl, "locale": get_document_locale(self.prs)}

    def embed_brand_fonts(self, path: Path) -> None:
        """Embed fonts from the document brand into the file."""
        brand = self.load_brand()
        if brand is not None:
            from slides_factory.layout.font_embed import embed_fonts_in_pptx

            fonts = brand.fonts.embeddable_fonts(brand)
            if fonts:
                embed_fonts_in_pptx(path, fonts)
