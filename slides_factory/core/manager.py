"""Slide lifecycle management — handles insertion, deletion, and metadata retrieval.

Classes:
    SlideManager — Encapsulates the logic for manipulating slides within a presentation.
"""

from __future__ import annotations

from pptx import Presentation
from pptx.slide import Slide


class SlideManager:
    """Manages the collection of slides in a PowerPoint presentation."""

    def __init__(self, prs: Presentation):
        self.prs = prs

    def delete_slide(self, index: int) -> None:
        """Remove a slide and clean up its relationships and package parts."""
        slide = self.prs.slides[index]
        slide_ids = self.prs.slides._sldIdLst
        r_id = slide_ids[index].rId
        self.prs.part.drop_rel(r_id)
        del slide_ids[index]
        partname = slide.part.partname
        parts = getattr(self.prs.part.package, "_parts", None)
        if parts is not None and partname in parts:
            del parts[partname]

    def insert_slide(self, layout, index: int) -> Slide:
        """Add a slide and move it to the requested index."""
        self.prs.slides.add_slide(layout)
        slide_ids = self.prs.slides._sldIdLst
        new_id = slide_ids[-1]
        slide_ids.remove(new_id)
        slide_ids.insert(index, new_id)
        return self.prs.slides[index]

    def clear_slide_shapes(self, slide: Slide) -> None:
        """Remove custom shapes and clear placeholder text before re-render."""
        sp_tree = slide.shapes._spTree
        for shape in list(slide.shapes):
            if shape.is_placeholder:
                if shape.has_text_frame:
                    shape.text_frame.clear()
            else:
                sp_tree.remove(shape._element)

    def remove_slide(self, index: int) -> None:
        """Delete a slide by index with bounds checking."""
        if 0 <= index < len(self.prs.slides):
            self.delete_slide(index)
        else:
            raise IndexError(f"Slide index {index} out of range (0-{len(self.prs.slides) - 1})")
