"""Raw grid slide lifecycle — create, mutate cells, and re-render via LayoutEngine.

Classes:
    GridSlideService — CRUD for ``$grid`` metadata slides (CLI builder path).
"""

from __future__ import annotations

from typing import Any

from pptx import Presentation

from slides_factory.app import SlideFactory
from slides_factory.core.engine import LayoutEngine
from slides_factory.core.manager import SlideManager
from slides_factory.layout_spec import Layout
from slides_factory.metadata import read_metadata, write_metadata

RAW_LAYOUT_ID = "$grid"


class _UnsetType:
    """Sentinel for unset optional parameters."""

    def __repr__(self) -> str:
        return "UNSET"


_UNSET = _UnsetType()


def normalize_cell(cell: dict[str, Any]) -> dict[str, Any]:
    """Normalize cell element metadata, dropping legacy utility-string style."""
    entry = dict(cell)
    element = dict(entry.get("element", {}))
    style = element.get("style")
    if isinstance(style, str):
        element.pop("style", None)
    elif not isinstance(style, dict):
        element["style"] = {}
    else:
        element["style"] = dict(style)
    entry["element"] = element
    return entry


def normalize_layout_dict(layout: dict[str, Any]) -> dict[str, Any]:
    """Normalize a raw layout dict before Pydantic validation."""
    spec = dict(layout)
    spec["cells"] = [normalize_cell(dict(cell)) for cell in spec.get("cells", [])]
    spec.setdefault("frame_style", {})
    return spec


def merge_frame_info(
    *,
    base: dict[str, Any] | None = None,
    updates: dict[str, Any] | None = None,
) -> dict[str, Any]:
    """Merge user frame-info updates onto an optional base dict."""
    info = dict(base or {})
    if updates:
        info.update(updates)
    return info


def validate_element_props(kind: str, props: dict[str, Any], app: SlideFactory) -> None:
    """Validate raw props against a registered element before re-rendering."""
    app.get_element(kind).validate_props(props or {})


def validate_element_style(kind: str, style: dict[str, Any] | None, app: SlideFactory) -> None:
    """Validate raw style JSON against a registered element before re-rendering."""
    app.get_element(kind).validate_style(style)


class GridSlideService:
    """Manage raw grid slides stored under ``$grid`` metadata."""

    def __init__(self, prs: Presentation, app: SlideFactory):
        self.prs = prs
        self.app = app
        self.engine = LayoutEngine(prs, app=app)
        self.mgr = SlideManager(prs)

    def require_grid_data(self, index: int) -> dict[str, Any]:
        """Return a mutable copy of a grid slide's stored spec, or raise."""
        if index < 0 or index >= len(self.prs.slides):
            raise IndexError(f"Slide index {index} out of range (0-{len(self.prs.slides) - 1})")
        meta = read_metadata(self.prs.slides[index])
        if not meta or meta.get("template_id") != RAW_LAYOUT_ID:
            raise ValueError(
                f"Slide {index} is not a raw grid slide; cells can only be edited on slides "
                "created with 'slide new'."
            )
        data = meta.get("data")
        spec: dict[str, Any] = dict(data) if isinstance(data, dict) else {}
        spec["cells"] = [normalize_cell(dict(cell)) for cell in spec.get("cells", [])]
        return spec

    def _render_validated(
        self,
        slide,
        validated: Layout,
        *,
        frame: str | None,
        rtl: bool | None,
        locale: str | None,
        stored_frame: str | None = None,
    ) -> tuple[str | None, bool, str]:
        prep = self.engine.prepare_render(
            frame=frame,
            rtl=rtl,
            locale=locale,
            stored_frame=stored_frame,
            frame_style=validated.frame_style,
        )
        self.engine.render_frame(
            slide,
            prep.frame_tpl,
            prep.ctx,
            prep.brand,
            validated.frame_info,
            validated.frame_style,
        )
        self.engine.render_grid(slide, validated, prep.ctx)
        return prep.frame_id, prep.rtl, prep.locale

    def add_layout_slide(
        self,
        layout: dict[str, Any],
        *,
        at: int | None = None,
        frame: str | None = None,
        rtl: bool | None = None,
        locale: str | None = None,
    ) -> dict[str, Any]:
        """Render a raw grid Layout onto a new slide and store metadata."""
        validated = Layout.model_validate(normalize_layout_dict(layout))
        prep = self.engine.prepare_render(
            frame=frame,
            rtl=rtl,
            locale=locale,
            frame_style=validated.frame_style,
        )
        self.engine.ensure_frame_allows_layout(prep.frame_tpl)

        pptx_layout = self.engine.resolve_blank_layout()
        if at is None:
            slide = self.prs.slides.add_slide(pptx_layout)
            index = len(self.prs.slides) - 1
        else:
            slide = self.mgr.insert_slide(pptx_layout, at)
            index = at

        frame_id, active_rtl, active_locale = self._render_validated(
            slide,
            validated,
            frame=frame,
            rtl=rtl,
            locale=locale,
        )
        payload = validated.model_dump(mode="json")
        write_metadata(slide, RAW_LAYOUT_ID, payload, frame_id=frame_id)

        result: dict[str, Any] = {
            "slide_index": index,
            "kind": "grid",
            "rtl": active_rtl,
            "locale": active_locale,
            "data": payload,
        }
        if frame_id is not None:
            result["frame_id"] = frame_id
        return result

    def rerender_layout(
        self,
        index: int,
        layout: dict[str, Any],
        *,
        frame: str | None = None,
        rtl: bool | None = None,
        locale: str | None = None,
    ) -> dict[str, Any]:
        """Re-render an existing raw grid slide in place from a mutated layout."""
        validated = Layout.model_validate(layout)
        existing_meta = read_metadata(self.prs.slides[index])
        stored_frame = existing_meta.get("frame_id") if existing_meta else None
        slide = self.prs.slides[index]
        self.mgr.clear_slide_shapes(slide)
        frame_id, active_rtl, active_locale = self._render_validated(
            slide,
            validated,
            frame=frame,
            rtl=rtl,
            locale=locale,
            stored_frame=stored_frame,
        )
        payload = validated.model_dump(mode="json")
        write_metadata(slide, RAW_LAYOUT_ID, payload, frame_id=frame_id)

        result: dict[str, Any] = {
            "slide_index": index,
            "kind": "grid",
            "rtl": active_rtl,
            "locale": active_locale,
            "data": payload,
        }
        if frame_id is not None:
            result["frame_id"] = frame_id
        return result

    def new_grid_slide(
        self,
        *,
        grid: str = "",
        frame: str | None = None,
        frame_info: dict[str, Any] | None = None,
        frame_style: dict[str, Any] | None = None,
        at: int | None = None,
        rtl: bool | None = None,
        locale: str | None = None,
    ) -> dict[str, Any]:
        """Create an empty grid slide ready for cell append operations."""
        data: dict[str, Any] = {
            "frame_info": merge_frame_info(updates=frame_info),
            "frame_style": dict(frame_style or {}),
            "grid": grid,
            "cells": [],
        }
        return self.add_layout_slide(data, at=at, frame=frame, rtl=rtl, locale=locale)

    def add_cell(
        self,
        index: int,
        *,
        kind: str,
        at: str = "",
        props: dict[str, Any] | None = None,
        style: dict[str, Any] | None = None,
    ) -> dict[str, Any]:
        """Append an element to a grid slide and re-render it in place."""
        spec = self.require_grid_data(index)
        props = props or {}
        style = style or {}
        validate_element_props(kind, props, self.app)
        validate_element_style(kind, style, self.app)
        spec["cells"].append({"at": at, "element": {"kind": kind, "props": props, "style": style}})
        result = self.rerender_layout(index, spec)
        result["cell_index"] = len(spec["cells"]) - 1
        return result

    def set_cell(
        self,
        index: int,
        cell: int,
        *,
        kind: Any = _UNSET,
        at: Any = _UNSET,
        props: Any = _UNSET,
        style: Any = _UNSET,
    ) -> dict[str, Any]:
        """Update one cell on a grid slide; only provided fields change."""
        spec = self.require_grid_data(index)
        cells = spec["cells"]
        if cell < 0 or cell >= len(cells):
            raise IndexError(f"Cell index {cell} out of range (0-{len(cells) - 1})")

        entry = dict(cells[cell])
        element = dict(entry.get("element", {}))
        if isinstance(element.get("style"), str):
            element.pop("style", None)
        if kind is not _UNSET:
            element["kind"] = kind
        if props is not _UNSET:
            element["props"] = props
        if style is not _UNSET:
            element["style"] = style
        if at is not _UNSET:
            entry["at"] = at

        validate_element_props(element.get("kind", ""), element.get("props") or {}, self.app)
        validate_element_style(element.get("kind") or "", element.get("style"), self.app)
        entry["element"] = element
        cells[cell] = entry

        result = self.rerender_layout(index, spec)
        result["cell_index"] = cell
        return result

    def remove_cell(self, index: int, cell: int) -> dict[str, Any]:
        """Remove one cell from a grid slide and re-render it in place."""
        spec = self.require_grid_data(index)
        cells = spec["cells"]
        if cell < 0 or cell >= len(cells):
            raise IndexError(f"Cell index {cell} out of range (0-{len(cells) - 1})")
        cells.pop(cell)
        result = self.rerender_layout(index, spec)
        result["removed_cell"] = cell
        return result

    def set_slide(
        self,
        index: int,
        *,
        grid: Any = _UNSET,
        frame: str | None = None,
        frame_info: dict[str, Any] | None = None,
        frame_style: dict[str, Any] | None = None,
        rtl: bool | None = None,
        locale: str | None = None,
    ) -> dict[str, Any]:
        """Update slide-level settings (grid classes, frame info, frame) in place."""
        spec = self.require_grid_data(index)
        if grid is not _UNSET:
            spec["grid"] = grid
        if frame_info is not None:
            spec["frame_info"] = merge_frame_info(
                base=spec.get("frame_info") if isinstance(spec.get("frame_info"), dict) else None,
                updates=frame_info,
            )
        if frame_style is not None:
            spec["frame_style"] = dict(frame_style)
        return self.rerender_layout(index, spec, frame=frame, rtl=rtl, locale=locale)
