"""Document-level operations on .pptx files (Facade).

This module provides the high-level API for manipulating slides, acting as a facade
that delegates work to core services.
"""

from __future__ import annotations

from pathlib import Path
from typing import TYPE_CHECKING, Any

from pptx import Presentation

from slides_factory.brand import BrandTheme
from slides_factory.core.engine import LayoutEngine
from slides_factory.core.grid import _UNSET, RAW_LAYOUT_ID, GridSlideService
from slides_factory.core.manager import SlideManager
from slides_factory.core.session import PresentationSession
from slides_factory.frame import get_frame
from slides_factory.metadata import read_metadata, write_metadata

if TYPE_CHECKING:
    from slides_factory.app import SlideFactory

PACKAGE_ROOT = Path(__file__).resolve().parents[1]
DEFAULT_THEME = PACKAGE_ROOT / "themes" / "default.pptx"


def ensure_default_theme() -> Path:
    """Create the default theme if it does not exist."""
    from pptx.util import Emu

    if DEFAULT_THEME.exists():
        existing = Presentation(str(DEFAULT_THEME))
        if int(existing.slide_width) != Emu(9144000):
            DEFAULT_THEME.unlink()
        else:
            return DEFAULT_THEME

    DEFAULT_THEME.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    # Use a temporary SlideManager to clear slides during theme creation
    mgr = SlideManager(prs)
    while len(prs.slides) > 0:
        mgr.delete_slide(0)

    prs.save(DEFAULT_THEME)
    return DEFAULT_THEME


def open_document(path: Path) -> Presentation:
    """Load an existing .pptx presentation from disk."""
    return PresentationSession.open(path).presentation


def load_document_brand(prs: Presentation) -> BrandTheme | None:
    """Load brand YAML stored on the document, if configured."""
    return PresentationSession(prs).load_brand()


def new_presentation(
    theme: Path | None = None,
    *,
    brand: Path | None = None,
    rtl: bool = False,
    locale: str = "en",
) -> Presentation:
    """Create an in-memory presentation shell."""
    session = PresentationSession.create_new(
        theme=theme,
        brand=brand,
        rtl=rtl,
        locale=locale,
        default_theme_provider=ensure_default_theme,
        delete_slide_callback=lambda prs, idx: SlideManager(prs).delete_slide(idx),
    )
    return session.presentation


def create_document(
    output: Path,
    theme: Path | None = None,
    *,
    brand: Path | None = None,
    rtl: bool = False,
    locale: str = "en",
) -> Presentation:
    """Create a new empty presentation and save it to output."""
    prs = new_presentation(theme=theme, brand=brand, rtl=rtl, locale=locale)
    session = PresentationSession(prs)
    session.save(output)
    return Presentation(str(output))


def update_document_rtl(
    prs: Presentation,
    *,
    rtl: bool,
    locale: str | None = None,
) -> dict[str, str | bool]:
    """Update document RTL/locale flags."""
    return PresentationSession(prs).update_rtl(rtl, locale)


def save_document(prs: Presentation, path: Path) -> None:
    """Persist a presentation to disk and embed fonts."""
    session = PresentationSession(prs)
    session.save(path)
    session.embed_brand_fonts(path)


def delete_slide(prs: Presentation, index: int) -> None:
    SlideManager(prs).delete_slide(index)


def insert_slide(prs: Presentation, layout, index: int):
    return SlideManager(prs).insert_slide(layout, index)


# --- Render Coordination (Delegating to LayoutEngine) ---


def add_slide(
    prs: Presentation,
    template_id: str,
    data: dict[str, Any],
    *,
    app: SlideFactory,
    at: int | None = None,
    frame: str | None = None,
    rtl: bool | None = None,
    locale: str | None = None,
) -> dict[str, Any]:
    template = app.get_template(template_id)
    validated = template.validate_data(data)
    layout = template.resolve_layout(prs)

    engine = LayoutEngine(prs, app=app)
    mgr = SlideManager(prs)

    prep = engine.prepare_render(
        frame=frame,
        rtl=rtl,
        locale=locale,
        template_default=type(template).default_frame,
    )

    if at is None:
        slide = prs.slides.add_slide(layout)
        index = len(prs.slides) - 1
    else:
        slide = mgr.insert_slide(layout, at)
        index = at

    engine.render_frame(
        slide,
        prep.frame_tpl,
        prep.ctx,
        prep.brand,
        # Use template's helper methods for chrome/style extraction
        template.frame_chrome(validated),
        template.frame_style_data(validated),
    )
    template.render(slide, validated, prep.ctx)
    write_metadata(
        slide,
        template_id,
        validated.model_dump(mode="json"),
        frame_id=prep.frame_id if prep.brand else None,
    )

    return {
        "slide_index": index,
        "template_id": template_id,
        "rtl": prep.rtl,
        "locale": prep.locale,
        "data": validated.model_dump(mode="json"),
        **({"frame_id": prep.frame_id} if prep.brand else {}),
    }


def add_frame_slide(
    prs: Presentation,
    frame_id: str,
    data: dict[str, Any],
    *,
    app: SlideFactory,
    at: int | None = None,
    rtl: bool | None = None,
    locale: str | None = None,
) -> dict[str, Any]:
    engine = LayoutEngine(prs, app=app)
    mgr = SlideManager(prs)
    frame_tpl = get_frame(app, frame_id)
    validated = frame_tpl.validate_info(data)

    prep = engine.prepare_render(frame=frame_id, rtl=rtl, locale=locale)

    pptx_layout = engine.resolve_blank_layout()
    if at is None:
        slide = prs.slides.add_slide(pptx_layout)
        index = len(prs.slides) - 1
    else:
        slide = mgr.insert_slide(pptx_layout, at)
        index = at

    engine.render_frame(slide, prep.frame_tpl, prep.ctx, prep.brand, validated)
    payload = validated.model_dump(mode="json")
    write_metadata(slide, "$frame", payload, frame_id=prep.frame_id)

    return {
        "slide_index": index,
        "kind": "frame",
        "frame_id": prep.frame_id,
        "rtl": prep.rtl,
        "locale": prep.locale,
        "data": payload,
    }


def edit_slide(
    prs: Presentation,
    index: int,
    data: dict[str, Any],
    *,
    app: SlideFactory,
    template_id: str | None = None,
    frame: str | None = None,
    rtl: bool | None = None,
    locale: str | None = None,
) -> dict[str, Any]:
    mgr = SlideManager(prs)
    if index < 0 or index >= len(prs.slides):
        raise IndexError(f"Slide index {index} out of range")

    existing_meta = read_metadata(prs.slides[index])
    res_tid = template_id or (existing_meta.get("template_id") if existing_meta else None)
    if not res_tid:
        raise ValueError("Slide has no template metadata. Pass --template explicitly.")

    template = app.get_template(res_tid)
    validated = template.validate_data(data)
    engine = LayoutEngine(prs, app=app)

    old_tid = existing_meta.get("template_id") if existing_meta else None
    changing = bool(template_id and old_tid and template_id != old_tid)
    stored_frame = None if changing else (existing_meta.get("frame_id") if existing_meta else None)

    prep = engine.prepare_render(
        frame=frame,
        rtl=rtl,
        locale=locale,
        stored_frame=stored_frame,
        template_default=type(template).default_frame,
    )

    if changing:
        mgr.delete_slide(index)
        return add_slide(
            prs,
            res_tid,
            data,
            app=app,
            at=index,
            frame=frame,
            rtl=prep.rtl,
            locale=prep.locale,
        )

    slide = prs.slides[index]
    mgr.clear_slide_shapes(slide)
    engine.render_frame(
        slide,
        prep.frame_tpl,
        prep.ctx,
        prep.brand,
        template.frame_chrome(validated),
        template.frame_style_data(validated),
    )
    template.render(slide, validated, prep.ctx)
    write_metadata(
        slide,
        res_tid,
        validated.model_dump(mode="json"),
        frame_id=prep.frame_id if prep.brand else None,
    )

    return {
        "slide_index": index,
        "template_id": res_tid,
        "rtl": prep.rtl,
        "locale": prep.locale,
        "data": validated.model_dump(mode="json"),
        **({"frame_id": prep.frame_id} if prep.brand else {}),
    }


def remove_slide(prs: Presentation, index: int) -> None:
    SlideManager(prs).remove_slide(index)


def _grid_service(prs: Presentation, app: SlideFactory) -> GridSlideService:
    return GridSlideService(prs, app=app)


def add_layout_slide(
    prs: Presentation,
    layout: dict[str, Any],
    *,
    app: SlideFactory,
    at: int | None = None,
    frame: str | None = None,
    rtl: bool | None = None,
    locale: str | None = None,
) -> dict[str, Any]:
    """Render a raw grid Layout (no template) onto a new slide and store it."""
    return _grid_service(prs, app=app).add_layout_slide(
        layout, at=at, frame=frame, rtl=rtl, locale=locale
    )


def new_grid_slide(
    prs: Presentation,
    *,
    app: SlideFactory,
    grid: str = "",
    frame: str | None = None,
    frame_info: dict[str, Any] | None = None,
    frame_style: dict[str, Any] | None = None,
    at: int | None = None,
    rtl: bool | None = None,
    locale: str | None = None,
) -> dict[str, Any]:
    """Create an empty grid slide ready for ``add_cell`` calls."""
    return _grid_service(prs, app=app).new_grid_slide(
        grid=grid,
        frame=frame,
        frame_info=frame_info,
        frame_style=frame_style,
        at=at,
        rtl=rtl,
        locale=locale,
    )


def add_cell(
    prs: Presentation,
    index: int,
    *,
    app: SlideFactory,
    kind: str,
    at: str = "",
    props: dict[str, Any] | None = None,
    style: dict[str, Any] | None = None,
) -> dict[str, Any]:
    """Append an element to a grid slide and re-render it in place."""
    return _grid_service(prs, app=app).add_cell(index, kind=kind, at=at, props=props, style=style)


def set_cell(
    prs: Presentation,
    index: int,
    cell: int,
    *,
    app: SlideFactory,
    kind: Any = _UNSET,
    at: Any = _UNSET,
    props: Any = _UNSET,
    style: Any = _UNSET,
) -> dict[str, Any]:
    """Update one cell on a grid slide; only provided fields change."""
    return _grid_service(prs, app=app).set_cell(
        index, cell, kind=kind, at=at, props=props, style=style
    )


def remove_cell(
    prs: Presentation,
    index: int,
    cell: int,
    *,
    app: SlideFactory,
) -> dict[str, Any]:
    """Remove one cell from a grid slide and re-render it in place."""
    return _grid_service(prs, app=app).remove_cell(index, cell)


def set_slide(
    prs: Presentation,
    index: int,
    *,
    app: SlideFactory,
    grid: Any = _UNSET,
    frame: str | None = None,
    frame_info: dict[str, Any] | None = None,
    frame_style: dict[str, Any] | None = None,
    rtl: bool | None = None,
    locale: str | None = None,
) -> dict[str, Any]:
    """Update slide-level settings (grid classes, frame info, frame) in place."""
    return _grid_service(prs, app=app).set_slide(
        index,
        grid=grid,
        frame=frame,
        frame_info=frame_info,
        frame_style=frame_style,
        rtl=rtl,
        locale=locale,
    )


def get_slide_info(
    prs: Presentation, index: int, app: SlideFactory
) -> dict[str, Any]:
    """Return template id and JSON data for one slide."""
    if index < 0 or index >= len(prs.slides):
        raise IndexError(f"Slide index {index} out of range (0-{len(prs.slides) - 1})")

    slide = prs.slides[index]
    meta = read_metadata(slide)
    if meta:
        template_id = meta["template_id"]
        if template_id == RAW_LAYOUT_ID:
            return {
                "slide_index": index,
                "template_id": None,
                "kind": "grid",
                "frame_id": meta.get("frame_id"),
                "data": meta["data"],
            }
        if template_id == "$frame":
            return {
                "slide_index": index,
                "template_id": None,
                "kind": "frame",
                "frame_id": meta.get("frame_id"),
                "data": meta["data"],
            }
        template = app.get_template(template_id)
        validated = template.validate_data(meta["data"])
        return {
            "slide_index": index,
            "template_id": template_id,
            "data": validated.model_dump(mode="json"),
            **({"frame_id": meta.get("frame_id")} if meta.get("frame_id") else {}),
        }

    return {
        "slide_index": index,
        "template_id": None,
        "data": {},
        "note": "No template metadata found",
    }


def list_slides_info(prs: Presentation, app: SlideFactory) -> dict[str, Any]:
    """Return summary of all slides."""
    slides = []
    for index in range(len(prs.slides)):
        meta = read_metadata(prs.slides[index])
        entry = {
            "index": index,
            "template_id": meta.get("template_id") if meta else None,
        }
        slides.append(entry)
    return {
        "slide_count": len(prs.slides),
        "slides": slides,
    }
