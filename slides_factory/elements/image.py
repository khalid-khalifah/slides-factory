"""Image element — renders a picture in a grid cell with fit modes and rounded corners."""

from __future__ import annotations

from pathlib import Path
from typing import Literal

from PIL import Image as PILImage
from pptx.enum.shapes import MSO_SHAPE
from pptx.slide import Slide
from pydantic import BaseModel, Field

from slides_factory.color_utils import hex_to_rgb
from slides_factory.layout.pct import image_aspect_ratio
from slides_factory.render_context import RenderContext
from slides_factory.styling import theme

_SUPPORTED_EXTENSIONS = frozenset({".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff", ".webp"})


class ImageProps(BaseModel):
    """Content props for the image element."""

    src: str = Field(description="Path to the image file (PNG, JPG, SVG, etc.).")
    alt: str = Field(default="", description="Accessible alt text for the image.")
    fit: Literal["contain", "cover", "stretch", "fill"] = Field(
        default="contain",
        description=(
            "How the image fits its cell: "
            "'contain' (fit inside, preserve ratio, letterbox), "
            "'cover' (fill cell, preserve ratio, crop overflow), "
            "'stretch' (fill cell exactly, non-uniform scale)"
        ),
    )


class ImageStyle(BaseModel):
    """Look overrides for the image element."""

    radius: str = Field(
        default="none",
        description="Theme radius token (e.g. 'none', 'sm', 'md', 'lg').",
    )
    opacity: float = Field(
        default=1.0,
        ge=0.0,
        le=1.0,
        description="Image opacity (0–1). Applied via shape fill when < 1.0.",
    )


def _resolve_image_path(src: str, ctx: RenderContext) -> Path:
    """Resolve ``src`` to an absolute ``Path``, checking brand logos and filesystem.

    Resolution order:
    1. If ``src`` matches a brand logo key, resolve from brand.
    2. Otherwise resolve relative to CWD or as absolute path.
    """
    path = Path(src)

    # If the path is explicitly absolute or relative and exists, use it directly.
    if path.is_absolute():
        if path.is_file():
            return path.resolve()
        raise FileNotFoundError(f"image not found at {path.resolve()}")
    if path.is_file():
        return path.resolve()

    # Check relative to working directory.
    cwd_path = Path.cwd() / path
    if cwd_path.is_file():
        return cwd_path.resolve()

    # Try resolving from brand logos if brand is available.
    if ctx.brand is not None:
        try:
            brand_path = ctx.brand.resolve_logo_raster(src)
            if brand_path is not None and brand_path.is_file():
                return brand_path.resolve()
        except (KeyError, ValueError, OSError, ImportError):
            pass

    raise FileNotFoundError(f"image not found: {src}")


def _fit_contain(
    box_w: int, box_h: int, img_w: int, img_h: int
) -> tuple[int, int, int, int]:
    """Return (left, top, width, height) for a contain-fitted image."""
    img_aspect = img_w / img_h
    box_aspect = box_w / box_h

    if img_aspect > box_aspect:
        # Image is wider: fit to width
        draw_w = box_w
        draw_h = int(box_w / img_aspect)
    else:
        # Image is taller: fit to height
        draw_h = box_h
        draw_w = int(box_h * img_aspect)

    left = (box_w - draw_w) // 2
    top = (box_h - draw_h) // 2
    return left, top, draw_w, draw_h


def _fit_cover(
    box_w: int, box_h: int, img_w: int, img_h: int
) -> tuple[int, int, int, int, float, float, float, float]:
    """Return (left, top, width, height, crop_left, crop_top, crop_right, crop_bottom).

    Crop values are python-pptx picture crop fractions (0–1) relative to the unscaled image.
    """
    img_aspect = img_w / img_h
    box_aspect = box_w / box_h

    if img_aspect > box_aspect:
        # Image is wider: fit to height, crop sides
        draw_h = box_h
        draw_w = int(box_h * img_aspect)
        overflow = draw_w - box_w
        crop_l = overflow / (2 * draw_w)
        crop_r = overflow / (2 * draw_w)
        crop_t = 0.0
        crop_b = 0.0
        left = 0
        top = 0
    else:
        # Image is taller: fit to width, crop top/bottom
        draw_w = box_w
        draw_h = int(box_w / img_aspect)
        overflow = draw_h - box_h
        crop_t = overflow / (2 * draw_h)
        crop_b = overflow / (2 * draw_h)
        crop_l = 0.0
        crop_r = 0.0
        left = 0
        top = 0

    # Clamp crop values to avoid < 1px visible region
    crop_l = min(crop_l, 0.499)
    crop_r = min(crop_r, 0.499)
    crop_t = min(crop_t, 0.499)
    crop_b = min(crop_b, 0.499)

    return left, top, draw_w, draw_h, crop_l, crop_t, crop_r, crop_b


def _add_picture_with_radius(
    slide: Slide,
    image_path: str | Path,
    left: int,
    top: int,
    width: int,
    height: int,
    radius_token: str,
):
    """Add a picture, optionally with rounded corners.

    When ``radius_token`` is ``"none"`` (the default), adds the picture directly.
    Otherwise modifies the picture's XML to set a rounded rectangle preset geometry.
    """
    from lxml import etree

    pic = slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)

    if radius_token != "none":
        radius_fraction = theme.radius(radius_token)
        ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"

        # spPr is in the presentationML namespace, not drawingML.
        ns_p = "http://schemas.openxmlformats.org/presentationml/2006/main"
        sp_pr = pic._element.find(f"{{{ns_p}}}spPr")
        if sp_pr is not None:
            prst_geom = sp_pr.find(f"{{{ns_a}}}prstGeom")
            if prst_geom is not None:
                prst_geom.set("prst", "roundRect")
                # Set the adjustment value for corner radius.
                av_lst = prst_geom.find(f"{{{ns_a}}}avLst")
                if av_lst is not None:
                    gd = etree.SubElement(av_lst, f"{{{ns_a}}}gd")
                    gd.set("name", "adj")
                    # Scale radius: theme.radius returns fraction, adjust for pptx adj scale.
                    adj_val = int(round(min(0.5, max(0.0, radius_fraction)) * 100000))
                    gd.set("fmla", f"val {adj_val}")

    return pic


def render_image(
    slide: Slide,
    box: tuple[int, int, int, int],
    props: ImageProps,
    style: ImageStyle,
    ctx: RenderContext,
):
    """Render an image into the given cell box.

    Supports four fit modes: ``contain``, ``cover``, ``stretch``, and ``fill``.
    """
    left, top, width, height = box

    # Resolve the image path.
    image_path = _resolve_image_path(props.src, ctx)

    # Validate format.
    suffix = image_path.suffix.lower()
    if suffix not in _SUPPORTED_EXTENSIONS and suffix != ".svg":
        supported = ", ".join(sorted(_SUPPORTED_EXTENSIONS))
        raise ValueError(
            f"unsupported image format '{suffix}' for {image_path.name}. "
            f"Supported: {supported}"
        )

    # Handle SVG source — convert to PNG before Pillow can open it.
    if suffix == ".svg":
        from slides_factory.brand.logos import resolve_raster_logo

        raster_path = resolve_raster_logo(image_path)
    else:
        raster_path = image_path

    # Open the image to get dimensions.
    try:
        with PILImage.open(raster_path) as img:
            img_w, img_h = img.size
    except Exception as exc:
        raise ValueError(f"failed to open image {raster_path}: {exc}") from exc

    if img_w <= 0 or img_h <= 0:
        raise ValueError(f"image has invalid dimensions ({img_w}x{img_h}): {image_path}")

    fit = props.fit

    if fit in ("stretch", "fill"):
        # Non-uniform scale to fill box exactly
        draw_left = 0
        draw_top = 0
        draw_w = width
        draw_h = height
        crop_l = crop_t = crop_r = crop_b = 0.0
    elif fit == "contain":
        draw_left, draw_top, draw_w, draw_h = _fit_contain(width, height, img_w, img_h)
        crop_l = crop_t = crop_r = crop_b = 0.0
        # Offset by cell position
        draw_left += left
        draw_top += top
    elif fit == "cover":
        cropped = _fit_cover(width, height, img_w, img_h)
        draw_left, draw_top, draw_w, draw_h, crop_l, crop_t, crop_r, crop_b = cropped
        draw_left += left
        draw_top += top
    else:
        raise ValueError(f"unknown fit mode: {fit!r}")

    # Add the picture.
    pic = _add_picture_with_radius(
        slide, str(raster_path), draw_left, draw_top, draw_w, draw_h,
        style.radius,
    )

    # Apply cover crop values.
    if fit == "cover":
        pic.crop_left = crop_l
        pic.crop_top = crop_t
        pic.crop_right = crop_r
        pic.crop_bottom = crop_b

    # Set alt text (via cNvPr descr attribute — python-pptx has no direct API).
    if props.alt:
        ns_p = "http://schemas.openxmlformats.org/presentationml/2006/main"
        nv_pic_pr = pic._element.find(f"{{{ns_p}}}nvPicPr")
        if nv_pic_pr is not None:
            c_nv_pr = nv_pic_pr.find(f"{{{ns_p}}}cNvPr")
            if c_nv_pr is not None:
                c_nv_pr.set("descr", props.alt)

    # Apply opacity if < 1.0 (via XML manipulation).
    if style.opacity < 1.0:
        _apply_opacity(pic, style.opacity)


def _apply_opacity(pic, opacity: float) -> None:
    """Set picture opacity via XML namespace manipulation.

    Modifies the ``<a:blip>`` element inside the picture's ``<p:blipFill>``
    to add a ``<a:alpha>`` fixed value modifier.
    """
    from lxml import etree

    ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"

    blip = pic._element.find(f".//{{{ns_a}}}blip")
    if blip is None:
        return

    # Remove any existing alpha modifier.
    for existing in blip.findall(f"{{{ns_a}}}alpha"):
        blip.remove(existing)

    # Add the alpha modifier.
    alpha_el = etree.SubElement(blip, f"{{{ns_a}}}alpha")
    alpha_val = int(round(opacity * 100000))
    alpha_val = max(0, min(100000, alpha_val))
    alpha_el.set("val", str(alpha_val))
