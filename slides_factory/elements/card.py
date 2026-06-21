"""Card element — a filled rounded box with stat text."""

from __future__ import annotations

from pydantic import BaseModel
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR
from pptx.slide import Slide

from slides_factory.color_utils import hex_to_rgb
from slides_factory.elements.base import Box, style_paragraph
from slides_factory.layout.fonts import apply_shape_font
from slides_factory.render_context import RenderContext
from slides_factory.styling import theme
from slides_factory.styling.models import CardStyle, is_brand_fill_ref


_PALETTE_TEXT_DEFAULTS = frozenset({"primary", "muted"})


def _text_color_for_card(style: CardStyle, *, role: str) -> str:
    """Pick text color, defaulting to brand contrast when background is a brand fill."""
    if is_brand_fill_ref(style.background_color):
        contrast_ref = f"on-{style.background_color}"
        token = getattr(style, f"{role}_color")
        if token in _PALETTE_TEXT_DEFAULTS:
            return contrast_ref
    return getattr(style, f"{role}_color")


class CardProps(BaseModel):
    """Content props for the card element."""

    title: str = ""
    value: str = ""
    body: str = ""


def render_card(
    slide: Slide,
    box: Box,
    props: CardProps,
    style: CardStyle,
    ctx: RenderContext,
) -> None:
    """Render a filled card with optional title / value / body text."""
    left, top, width, height = box
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )

    if shape.adjustments:
        shape.adjustments[0] = min(0.5, max(0.0, theme.radius(style.radius)))

    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(
        theme.resolve_style_color(style.background_color, ctx)
    )
    shape.line.fill.background()

    frame = shape.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    rows: list[tuple[str, float, bool, str]] = []
    if props.title:
        rows.append(
            (
                props.title,
                theme.font_size_pt(style.title_size),
                False,
                _text_color_for_card(style, role="title"),
            )
        )
    if props.value:
        rows.append(
            (
                props.value,
                theme.font_size_pt(style.value_size),
                style.value_bold,
                _text_color_for_card(style, role="value"),
            )
        )
    if props.body:
        rows.append(
            (
                props.body,
                theme.font_size_pt(style.body_size),
                False,
                _text_color_for_card(style, role="body"),
            )
        )
    if not rows:
        rows.append(
            (
                "",
                theme.font_size_pt(style.body_size),
                False,
                _text_color_for_card(style, role="body"),
            )
        )

    for index, (content, size_pt, bold, color_token) in enumerate(rows):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = content
        style_paragraph(
            paragraph,
            ctx,
            size_pt=size_pt,
            bold=bold,
            color_token=color_token,
            align="center",
        )

    apply_shape_font(shape, ctx, style.font)
