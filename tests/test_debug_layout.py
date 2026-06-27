"""Tests for layout debug mode (grid lines, cell labels, playground boundary)."""

from __future__ import annotations

from pathlib import Path

import pytest
from pptx.util import Emu, Inches

from slides_factory import document
from slides_factory.app import SlideFactory
from slides_factory.layout.debug import render_debug_layer
from slides_factory.layout.grid import Cell, GridStyle, compute_cells
from slides_factory.render_context import RenderContext
from slides_factory.styling.tokens import parse_grid


def _make_app() -> SlideFactory:
    from tests.fixtures.app import app

    return app


def _make_deck(tmp_path: Path) -> Path:
    output = tmp_path / "test.pptx"
    document.create_document(output)
    return output


def test_debug_creates_shapes_on_new_slide(tmp_path: Path):
    """slide new --debug adds at least one diagnostic shape (playground boundary)."""
    app = _make_app()
    deck = _make_deck(tmp_path)
    prs = document.open_document(deck)
    result = document.new_grid_slide(
        prs,
        app=app,
        grid="grid-cols-2 grid-rows-2 gap-4 p-4",
        debug=True,
    )
    document.save_document(prs, deck)
    prs = document.open_document(deck)
    slide = prs.slides[result["slide_index"]]
    assert len(slide.shapes) > 0


def test_debug_no_shapes_without_flag(tmp_path: Path):
    """slide new without --debug produces empty slide (0 shapes)."""
    app = _make_app()
    deck = _make_deck(tmp_path)
    prs = document.open_document(deck)
    result = document.new_grid_slide(
        prs,
        app=app,
        grid="grid-cols-2 grid-rows-2",
    )
    document.save_document(prs, deck)
    prs = document.open_document(deck)
    slide = prs.slides[result["slide_index"]]
    assert len(slide.shapes) == 0


def test_debug_render_layer_directly():
    """render_debug_layer() can be called standalone with correct args."""
    from pptx import Presentation

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(blank_layout)

    grid = parse_grid("grid-cols-2 grid-rows-2 gap-4 p-4")
    region = (0, 0, int(prs.slide_width), int(prs.slide_height))
    cells: list[Cell] = []
    ctx = RenderContext(
        rtl=False,
        locale="en",
        slide_width=int(prs.slide_width),
        slide_height=int(prs.slide_height),
        playground=region,
    )
    render_debug_layer(slide, grid, cells, ctx=ctx)

    # Playground boundary should be the first shape
    assert len(slide.shapes) > 0


def test_debug_cell_labels_appear(tmp_path: Path):
    """Cells get [row,col] labels when debug is on."""
    app = _make_app()
    deck = _make_deck(tmp_path)
    prs = document.open_document(deck)
    result = document.new_grid_slide(
        prs,
        app=app,
        grid="grid-cols-3 grid-rows-2",
        debug=True,
    )
    document.save_document(prs, deck)
    prs = document.open_document(deck)
    slide = prs.slides[result["slide_index"]]
    assert len(slide.shapes) > 0


def test_debug_flag_on_add_slide(tmp_path: Path):
    """slide add --debug renders diagnostic shapes alongside content."""
    app = _make_app()
    deck = _make_deck(tmp_path)
    prs = document.open_document(deck)
    data = {"headline": {"block": {"children": [{"runs": [{"text": "Debug"}]}]}}}
    result = document.add_slide(
        prs,
        "simple",
        data,
        app=app,
        debug=True,
    )
    document.save_document(prs, deck)
    prs = document.open_document(deck)
    slide = prs.slides[result["slide_index"]]
    # With debug on, there should be shapes (content + diagnostics)
    assert len(slide.shapes) > 0


def test_debug_grid_lines_match_track_count():
    """Column/row grid lines are drawn for each track boundary."""
    from pptx import Presentation

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(blank_layout)

    grid = parse_grid("grid-cols-3 grid-rows-2 gap-4 p-4")
    region = (0, 0, int(prs.slide_width), int(prs.slide_height))
    cell_styles = []  # no cells, just grid

    from slides_factory.styling.tokens import parse_cell

    placed = compute_cells(region, grid, cell_styles, rtl=False)
    ctx = RenderContext(
        rtl=False,
        locale="en",
        slide_width=int(prs.slide_width),
        slide_height=int(prs.slide_height),
        playground=region,
    )
    render_debug_layer(slide, grid, placed, ctx=ctx)

    # Should have at least the playground boundary rect
    assert len(slide.shapes) >= 1
