"""Frame with a declared playground region and an information layer.

Exercises the frame signature with a frame-specific info model: title band and
optional page number in the footer.
"""

from __future__ import annotations

from pptx.slide import Slide
from pptx.util import Emu
from pydantic import BaseModel, Field

from slides_factory.layout.pct import PctBox
from slides_factory.render_context import RenderContext
from tests.fixtures.app import app
from tests.fixtures.palettes import TEST_LIGHT


class PaneledInfo(BaseModel):
    title: str | None = Field(default=None, description="Title band text.")
    page_number: int | None = Field(default=None, ge=0, description="Footer page number.")


@app.frame(
    "paneled",
    name="Paneled",
    description="Frame with playground + info layer for grid tests",
    palette=TEST_LIGHT,
    playground=PctBox(left=10, top=25, width=80, height=65),
    frame_info_model=PaneledInfo,
)
def paneled(slide: Slide, ctx: RenderContext, info: PaneledInfo) -> None:
    if info.title:
        textbox = slide.shapes.add_textbox(
            Emu(0), Emu(0), Emu(ctx.slide_width), Emu(int(ctx.slide_height * 0.15))
        )
        textbox.text_frame.text = info.title
    if info.page_number is not None:
        footer = slide.shapes.add_textbox(
            Emu(0),
            Emu(int(ctx.slide_height * 0.92)),
            Emu(ctx.slide_width),
            Emu(int(ctx.slide_height * 0.08)),
        )
        footer.text_frame.text = str(info.page_number)
