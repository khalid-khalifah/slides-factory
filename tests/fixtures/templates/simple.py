"""Minimal title-and-content template for core document tests."""

from typing import Annotated

from pydantic import Field
from pptx.slide import Slide

from slides_factory.render_context import RenderContext
from slides_factory.template_input import TemplateInput
from tests.fixtures.app import app


class SimpleInput(TemplateInput):
    title: Annotated[str, Field(description="Slide title")]
    body: Annotated[str, Field(description="Body text")] = ""


def _extract_simple(slide: Slide):
    title = slide.shapes.title.text if slide.shapes.title else ""
    body = ""
    if len(slide.placeholders) > 1:
        body = slide.placeholders[1].text
    return {"title": title, "body": body}


@app.template(
    "simple",
    name="Simple",
    description="Title and optional body for core tests",
    layout_name="Title and Content",
    tags=["content", "test"],
    extract=_extract_simple,
)
def simple(slide: Slide, ctx: RenderContext, data: SimpleInput) -> None:
    if slide.shapes.title:
        slide.shapes.title.text = data.title
    if data.body and len(slide.placeholders) > 1:
        slide.placeholders[1].text = data.body
