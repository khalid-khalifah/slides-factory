"""Template with a default frame for resolution tests."""

from typing import Annotated

from pydantic import Field
from pptx.slide import Slide

from slides_factory.render_context import RenderContext
from slides_factory.template_input import TemplateInput
from tests.fixtures.app import app


class PairedInput(TemplateInput):
    title: Annotated[str, Field(description="Slide title")]


def _extract_paired(slide: Slide):
    title = slide.shapes.title.text if slide.shapes.title else ""
    return {"title": title}


@app.template(
    "paired",
    name="Paired",
    description="Template with default_frame=plain",
    layout_name="Title and Content",
    default_frame="plain",
    extract=_extract_paired,
)
def paired(slide: Slide, ctx: RenderContext, data: PairedInput) -> None:
    if slide.shapes.title:
        slide.shapes.title.text = data.title
