"""Template with required fields only — for validation edge-case tests."""

from typing import Annotated

from pydantic import Field
from pptx.slide import Slide

from slides_factory.render_context import RenderContext
from slides_factory.template_input import TemplateInput
from tests.fixtures.app import app


class StrictInput(TemplateInput):
    title: Annotated[str, Field(description="Required title")]
    count: Annotated[int, Field(description="Required count")]


def _extract_strict(slide: Slide):
    if slide.shapes.title is None:
        raise ValueError("strict slide missing title")
    text = slide.shapes.title.text
    title, _, rest = text.partition(" (")
    if not rest.endswith(")"):
        raise ValueError(f"cannot parse strict slide title: {text!r}")
    return {"title": title, "count": int(rest[:-1])}


@app.template(
    "strict",
    name="Strict",
    description="All fields required",
    layout_name="Title and Content",
    tags=["test"],
    extract=_extract_strict,
)
def strict(slide: Slide, ctx: RenderContext, data: StrictInput) -> None:
    if slide.shapes.title:
        slide.shapes.title.text = f"{data.title} ({data.count})"
