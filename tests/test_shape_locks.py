"""Tests for OOXML shape locking."""

from __future__ import annotations

import zipfile
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.util import Inches

from slides_factory.layout.locks import lock_shape, lock_shapes_added

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
NSMAP = {"a": A_NS}


def _make_dot_png(path: Path) -> Path:
    from PIL import Image

    Image.new("RGB", (1, 1), "red").save(path)
    return path


def _locks_xml(path: Path) -> list[tuple[str, dict[str, str]]]:
    with zipfile.ZipFile(path) as zf:
        slide = etree.fromstring(zf.read("ppt/slides/slide1.xml"))
    out: list[tuple[str, dict[str, str]]] = []
    for tag, xpath in (
        ("picLocks", "//a:picLocks"),
        ("spLocks", "//a:spLocks"),
    ):
        for node in slide.xpath(xpath, namespaces=NSMAP):
            out.append((tag, dict(node.attrib)))
    return out


def test_lock_shape_picture(tmp_path: Path) -> None:
    out = tmp_path / "pic.pptx"
    pic_path = _make_dot_png(tmp_path / "dot.png")
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(pic_path), Inches(1), Inches(1), width=Inches(2))
    lock_shape(slide.shapes[-1])
    prs.save(str(out))

    locks = _locks_xml(out)
    pic_locks = [attrs for tag, attrs in locks if tag == "picLocks"]
    assert pic_locks
    assert pic_locks[0]["noMove"] == "1"
    assert pic_locks[0]["noResize"] == "1"
    assert pic_locks[0]["noRot"] == "1"
    assert "noTextEdit" not in pic_locks[0]


def test_lock_shape_textbox(tmp_path: Path) -> None:
    out = tmp_path / "text.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1)).text = "hi"
    lock_shape(slide.shapes[-1])
    prs.save(str(out))

    locks = _locks_xml(out)
    assert any(tag == "spLocks" and attrs.get("noTextEdit") == "1" for tag, attrs in locks)


def test_lock_shape_sets_user_drawn(tmp_path: Path) -> None:
    out = tmp_path / "userdrawn.pptx"
    pic_path = _make_dot_png(tmp_path / "dot.png")
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(pic_path), Inches(1), Inches(1), width=Inches(2))
    lock_shape(slide.shapes[-1])
    prs.save(str(out))

    P = "http://schemas.openxmlformats.org/presentationml/2006/main"
    with zipfile.ZipFile(out) as zf:
        slide_xml = etree.fromstring(zf.read("ppt/slides/slide1.xml"))
    nv_pr = slide_xml.xpath("//p:nvPicPr/p:nvPr", namespaces={"p": P})
    assert nv_pr and nv_pr[0].get("userDrawn") == "1"


def test_lock_shapes_added_only_new_shapes(tmp_path: Path) -> None:
    out = tmp_path / "added.pptx"
    pic_path = _make_dot_png(tmp_path / "dot.png")
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    existing = {id(s._element) for s in slide.shapes}
    slide.shapes.add_picture(str(pic_path), Inches(1), Inches(1), width=Inches(1))
    lock_shapes_added(slide, existing)
    prs.save(str(out))

    locks = _locks_xml(out)
    assert any(tag == "picLocks" for tag, _ in locks)
