"""Class-based templates built on top of the grid+element core.

Demonstrates inferred input models from @at cells plus template chrome fields.
The template is registered on a local factory so it does not pollute the
shared test catalog.
"""

from __future__ import annotations

import json
from pathlib import Path

import pytest
from pptx.enum.shapes import MSO_SHAPE_TYPE
from typer.testing import CliRunner

from slides_factory import document
from slides_factory.app import SlideFactory
from slides_factory.templating import Template, at

runner = CliRunner()


def _make_factory() -> SlideFactory:
    factory = SlideFactory("kpi-app")

    # Register a minimal text element so templates can use @at(kind="text").
    from contextlib import suppress

    from pydantic import BaseModel

    from slides_factory.converters.text import TextBlock, render_text_block
    from slides_factory.elements.base import Box, element_from_function
    from slides_factory.render_context import RenderContext
    from slides_factory.styling import theme

    class _TxtProps(BaseModel):
        block: dict | None = None

    class _TxtStyle(BaseModel):
        text_size: str = "base"
        text_color: str = "primary"
        bold: bool = False
        align: str = "left"
        font: str = "body"

    def _txt_render(slide, box: Box, props, style, ctx: RenderContext):
        block = TextBlock(children=[])
        if props.block and isinstance(props.block, dict):
            with suppress(Exception):
                block = TextBlock.model_validate(props.block)
        textbox = slide.shapes.add_textbox(*box)
        render_text_block(
            block, textbox.text_frame, ctx,
            base_size_pt=theme.font_size_pt(style.text_size),
            base_color=style.text_color,
            base_bold=style.bold,
            alignment=style.align,
            font_slot=style.font,
            vertical_anchor="top",
        )

    factory._elements["text"] = element_from_function(
        _txt_render, kind="text", props_model=_TxtProps, style_model=_TxtStyle,
    )

    @factory.template(
        "kpi-duo",
        name="KPI Duo",
        description="A bold heading over two KPI cards side by side.",
        grid="grid-cols-2 grid-rows-[1_2] gap-4",
    )
    class KpiDuo(Template):
        @at("col-span-2", kind="text")
        def heading(self): ...

        @at("", kind="text")
        def revenue(self): ...

        @at("", kind="text")
        def customers(self): ...

    return factory


@pytest.fixture
def kpi_factory():
    factory = _make_factory()
    yield factory


def _kpi_data() -> dict:
    return {
        "title": "Q3",

        "styles": {},
        "frame_style": {},
        "heading": {"block": {"children": [{"runs": [{"text": "Q3"}]}]}},
        "revenue": {"block": {"children": [{"runs": [{"text": "Revenue: $1.2M"}]}]}},
        "customers": {"block": {"children": [{"runs": [{"text": "Customers: 8,400"}]}]}},
    }


def test_build_maps_typed_data_to_layout(kpi_factory: SlideFactory):
    template = kpi_factory.get_template("kpi-duo")
    data = template.validate_data(_kpi_data())
    layout = template.build(data)

    assert layout.grid == "grid-cols-2 grid-rows-[1_2] gap-4"
    assert layout.frame_info["title"] == "Q3"
    kinds = [c.element.kind for c in layout.cells]
    assert kinds == ["text", "text", "text"]
    assert layout.cells[0].at == "col-span-2"
    assert layout.cells[0].element.props == {"block": {"children": [{"runs": [{"text": "Q3"}]}]}}
    assert layout.cells[0].element.style == {}
    assert layout.cells[1].element.props == {
        "block": {"children": [{"runs": [{"text": "Revenue: $1.2M"}]}]},
    }

def test_template_build_applies_cell_styles(kpi_factory: SlideFactory):
    template = kpi_factory.get_template("kpi-duo")
    data = template.validate_data(
        {
            **_kpi_data(),
            "styles": {"heading": {"text_size": "2xl", "font": "title"}},
        }
    )
    layout = template.build(data)
    assert layout.cells[0].element.style == {"text_size": "2xl", "font": "title"}


def test_template_round_trips_typed_input(kpi_factory: SlideFactory, tmp_path: Path):
    output = tmp_path / "kpi.pptx"
    prs = document.create_document(output)
    data = _kpi_data()
    result = document.add_slide(prs, "kpi-duo", data, app=kpi_factory)
    document.save_document(prs, output)

    assert result["template_id"] == "kpi-duo"

    prs = document.open_document(output)
    info = document.get_slide_info(prs, 0, app=kpi_factory)
    assert info["template_id"] == "kpi-duo"
    assert info["data"] == data

    slide = prs.slides[0]
    textboxes = [s for s in slide.shapes if s.has_text_frame and not s.is_placeholder]
    assert len(textboxes) >= 1


def test_template_requires_at_least_one_cell(kpi_factory: SlideFactory):
    factory = kpi_factory

    with pytest.raises(TypeError, match="at least one @at method"):

        @factory.template("empty", name="Empty", description="")
        class Empty(Template):
            pass


def test_template_rejects_frame_info_name_collision(kpi_factory: SlideFactory):
    factory = kpi_factory

    with pytest.raises(TypeError, match="conflicts with a template chrome field"):

        @factory.template("bad-name", name="Bad", description="", grid="grid-cols-1")
        class BadName(Template):
            @at("", kind="text")
            def title(self): ...


def test_cli_lists_and_inspects_templates(kpi_factory: SlideFactory):
    cli = kpi_factory.cli

    listing = json.loads(runner.invoke(cli, ["templates", "list", "--json"]).output)
    ids = {t["id"] for t in listing["data"]["templates"]}
    assert "kpi-duo" in ids

    inspect = json.loads(runner.invoke(cli, ["templates", "inspect", "kpi-duo", "--json"]).output)
    props = inspect["data"]["json_schema"]["properties"]
    assert {"title", "heading", "revenue", "customers", "styles", "frame_style"} <= set(props)
    assert "page_number" not in props
    assert "total_pages" not in props


def test_cli_slide_add_from_template(kpi_factory: SlideFactory, tmp_path: Path):
    cli = kpi_factory.cli
    deck = tmp_path / "deck.pptx"
    runner.invoke(cli, ["doc", "create", "-o", str(deck)])

    data = _kpi_data()
    result = runner.invoke(
        cli,
        [
            "slide", "add-many", str(deck),
            "--template", "kpi-duo",
            "--rows", json.dumps([data]),
            "--json",
        ],
    )
    payload = json.loads(result.output)
    assert payload["ok"] is True
    assert payload["data"]["results"][0]["template_id"] == "kpi-duo"

    got = json.loads(runner.invoke(cli, ["doc", "get", str(deck), "--index", "0", "--json"]).output)
    assert got["data"]["template_id"] == "kpi-duo"
    assert got["data"]["data"] == data
