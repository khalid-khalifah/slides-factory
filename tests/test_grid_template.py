"""End-to-end tests for the core grid layout via the document pipeline."""

from __future__ import annotations

from pathlib import Path

from pptx.enum.shapes import MSO_SHAPE_TYPE

from slides_factory import document


def _spec() -> dict:
    return {
        "grid": "grid-cols-2 gap-4",
        "cells": [
            {
                "element": {
                    "kind": "text",
                    "style": {"text_size": "lg", "bold": True},
                    "props": {"text": "Highlights"},
                }
            },
            {
                "element": {
                    "kind": "card",
                    "style": {"background_color": "surface"},
                    "props": {"title": "Revenue", "value": "$1.2M"},
                }
            },
        ],
    }


def test_layout_slide_round_trip_without_brand(tmp_path: Path):
    output = tmp_path / "grid.pptx"
    prs = document.create_document(output)
    result = document.add_layout_slide(prs, _spec())
    document.save_document(prs, output)

    assert result["kind"] == "grid"

    prs = document.open_document(output)
    info = document.get_slide_info(prs, 0)
    assert info["kind"] == "grid"
    assert info["template_id"] is None
    assert info["data"]["grid"] == "grid-cols-2 gap-4"
    assert info["data"]["cells"][0]["element"]["props"]["text"] == "Highlights"

    slide = prs.slides[0]
    has_textbox = any(s.has_text_frame and not s.is_placeholder for s in slide.shapes)
    has_card = any(s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE for s in slide.shapes)
    assert has_textbox and has_card


def test_layout_slide_with_frame_draws_info_layer(tmp_path: Path, minimal_brand_yaml: Path):
    output = tmp_path / "branded.pptx"
    prs = document.create_document(output, brand=minimal_brand_yaml)
    spec = _spec()
    spec["frame_info"] = {"title": "Quarterly Review", "page_number": 1}
    result = document.add_layout_slide(prs, spec, frame="paneled")
    document.save_document(prs, output)

    assert result["frame_id"] == "paneled"

    prs = document.open_document(output)
    slide = prs.slides[0]
    texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
    assert "Quarterly Review" in texts


def test_layout_cells_land_inside_frame_playground(tmp_path: Path, minimal_brand_yaml: Path):
    output = tmp_path / "branded.pptx"
    prs = document.create_document(output, brand=minimal_brand_yaml)
    spec = {
        "frame_info": {"title": "T"},
        "grid": "grid-cols-1",
        "cells": [
            {
                "element": {
                    "kind": "card",
                    "style": {"background_color": "surface"},
                    "props": {"value": "X"},
                }
            }
        ],
    }
    document.add_layout_slide(prs, spec, frame="paneled")

    slide = prs.slides[0]
    card = next(s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE)
    # paneled playground: left 10%, top 25%, width 80%, height 65% of 10in x 7.5in
    width = int(prs.slide_width)
    height = int(prs.slide_height)
    assert abs(card.left - int(width * 0.10)) < width * 0.02
    assert abs(card.top - int(height * 0.25)) < height * 0.02
