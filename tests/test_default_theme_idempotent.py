"""Tests for document-level defaults and metadata robustness."""

import warnings
from unittest.mock import MagicMock, patch


class TestDefaultThemeIdempotent:
    """Test ensure_default_theme() idempotency."""

    def test_default_theme_idempotent(self, tmp_path, monkeypatch):
        """Creating the default theme multiple times should not fail."""
        from slides_factory.document import ensure_default_theme

        # Point DEFAULT_THEME at a temp location for isolation
        monkeypatch.setattr("slides_factory.document.DEFAULT_THEME", tmp_path / "default.pptx")

        path1 = ensure_default_theme()
        path2 = ensure_default_theme()

        assert path1 == path2
        assert path1.exists()

    def test_default_theme_returns_existing(self, tmp_path, monkeypatch):
        """If theme already exists with correct size, return it without recreating."""
        from pptx.util import Emu

        from slides_factory.document import ensure_default_theme

        monkeypatch.setattr("slides_factory.document.DEFAULT_THEME", tmp_path / "default.pptx")

        # Create a valid theme file first
        ensure_default_theme()

        with patch("slides_factory.document.Presentation") as mock_pptx:
            mock_instance = MagicMock()
            mock_instance.slide_width = Emu(9144000)
            mock_pptx.return_value = mock_instance

            result = ensure_default_theme()

            assert result == tmp_path / "default.pptx"
            # Presentation is called once to verify file size, not to recreate
            assert mock_pptx.call_count == 1


class TestMetadataCorrupted:
    """Test metadata handling with corrupted data."""

    def test_read_metadata_corrupt_json(self, tmp_path, caplog):
        """Corrupted JSON in speaker notes should return None and emit a warning."""
        import logging

        from pptx import Presentation
        from pptx.util import Emu

        from slides_factory.metadata import read_metadata

        prs = Presentation()
        prs.slide_width = Emu(9144000)
        prs.slide_height = Emu(5400000)

        slide = prs.slides.add_slide(prs.slide_layouts[0])
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = "{invalid json!!!"

        with caplog.at_level(logging.WARNING):
            result = read_metadata(slide)

        assert result is None
        assert len(caplog.records) == 1
        assert "Corrupted slide metadata" in caplog.text
        assert "slide" in caplog.text.lower()

    def test_read_metadata_no_notes(self, tmp_path):
        """Slides without notes should return None without warning."""
        from pptx import Presentation
        from pptx.util import Emu

        from slides_factory.metadata import read_metadata

        prs = Presentation()
        prs.slide_width = Emu(9144000)
        prs.slide_height = Emu(5400000)

        slide = prs.slides.add_slide(prs.slide_layouts[0])
        # Remove notes if they exist
        if slide.has_notes_slide:
            slide.notes_slide.notes_text_frame.clear()

        with warnings.catch_warnings(record=True) as w:
            warnings.simplefilter("always")
            result = read_metadata(slide)

        assert result is None
        assert len(w) == 0

    def test_read_metadata_empty_notes(self, tmp_path):
        """Slides with empty notes should return None without warning."""
        from pptx import Presentation
        from pptx.util import Emu

        from slides_factory.metadata import read_metadata

        prs = Presentation()
        prs.slide_width = Emu(9144000)
        prs.slide_height = Emu(5400000)

        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.notes_slide.notes_text_frame.text = "   "

        with warnings.catch_warnings(record=True) as w:
            warnings.simplefilter("always")
            result = read_metadata(slide)

        assert result is None
        assert len(w) == 0
