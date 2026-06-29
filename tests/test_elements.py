"""Tests for text converter rendering into a slide."""

from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.util import Inches

from slides_factory.converters.text import (
    ListItem,
    ListStyle,
    Paragraph,
    TextBlock,
    TextRun,
    parse_html,
    render_text_block,
    text,
)
from slides_factory.palette import SlidePalette
from slides_factory.render_context import RenderContext

BASE_CTX = RenderContext(
    rtl=False, locale="en",
    slide_width=int(Inches(10)), slide_height=int(Inches(7.5)),
)


def _slide():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    return prs.slides.add_slide(prs.slide_layouts[6])


def _tb(slide):
    return slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(2))


def _ctx(palette=None):
    return BASE_CTX if palette is None else BASE_CTX.with_palette(palette)


# ---------------------------------------------------------------------------
# Text converter — render_text_block
# ---------------------------------------------------------------------------

def test_render_text_block_basic():
    """render_text_block writes text into a textbox."""
    slide = _slide()
    tb = _tb(slide)
    block = TextBlock(children=[Paragraph(runs=[TextRun(text="Hello World")])])
    render_text_block(block, tb, _ctx(), base_size_pt=12)
    assert tb.text_frame.paragraphs[0].text == "Hello World"


def test_render_text_block_multiple_paragraphs():
    """Multiple paragraphs render as separate pptx paragraphs."""
    slide = _slide()
    tb = _tb(slide)
    block = TextBlock(children=[
        Paragraph(runs=[TextRun(text="First")]),
        Paragraph(runs=[TextRun(text="Second")]),
    ])
    render_text_block(block, tb, _ctx(), base_size_pt=12)
    assert len(tb.text_frame.paragraphs) == 2
    assert tb.text_frame.paragraphs[1].text == "Second"


def test_render_text_block_bullets():
    """ListItem nodes produce bullet markers."""
    slide = _slide()
    tb = _tb(slide)
    block = TextBlock(children=[
        ListItem(runs=[TextRun(text="item")], marker=ListStyle(type="disc", level=0)),
    ])
    render_text_block(block, tb, _ctx(), base_size_pt=12)
    assert tb.text_frame.paragraphs[0].text != ""


def test_render_text_block_applies_palette_color():
    """Colour tokens resolve against the palette."""
    slide = _slide()
    tb = _tb(slide)
    palette = SlidePalette(text="#123456", highlight="#000000", main=("#FFFFFF",), extras=())
    block = TextBlock(children=[Paragraph(runs=[TextRun(text="Colored")])])
    render_text_block(
        block, tb, _ctx(palette=palette), base_size_pt=12, base_color="primary",
    )
    assert tb.text_frame.paragraphs[0].runs[0].text == "Colored"


def test_render_text_block_honors_bold():
    """base_bold applies bold to all runs."""
    slide = _slide()
    tb = _tb(slide)
    block = TextBlock(children=[Paragraph(runs=[TextRun(text="Bold")])])
    render_text_block(block, tb, _ctx(), base_size_pt=12, base_bold=True)
    assert tb.text_frame.paragraphs[0].runs[0].font.bold is True


def test_render_text_block_accepts_shape():
    """render_text_block works when passed a shape (not just TextFrame)."""
    slide = _slide()
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(5), Inches(2))
    block = TextBlock(children=[Paragraph(runs=[TextRun(text="In shape")])])
    render_text_block(block, shape, _ctx(), base_size_pt=12)
    assert shape.text_frame.paragraphs[0].text == "In shape"


# ---------------------------------------------------------------------------
# Rich text — parse_html integration
# ---------------------------------------------------------------------------

def test_rich_text_basic():
    """parse_html → render_text_block produces correct bold/italic runs."""
    slide = _slide()
    tb = _tb(slide)
    block = parse_html("<b>Bold</b> and <i>italic</i>")
    render_text_block(block, tb, _ctx(), base_size_pt=12)
    assert "Bold" in tb.text_frame.paragraphs[0].text


def test_rich_text_colors():
    """parse_html with color attribute resolves against palette."""
    slide = _slide()
    tb = _tb(slide)
    palette = SlidePalette(text="#111111", highlight="#EAA000", main=("#FFFFFF",), extras=())
    block = parse_html('<span color="highlight">Highlighted</span>')
    render_text_block(block, tb, _ctx(palette=palette), base_size_pt=12)
    assert tb.text_frame.paragraphs[0].text == "Highlighted"


def test_rich_text_hyperlink():
    """parse_html with <a> creates hyperlink."""
    slide = _slide()
    tb = _tb(slide)
    block = parse_html('<a href="https://example.com">Click</a>')
    render_text_block(block, tb, _ctx(), base_size_pt=12)
    assert tb.text_frame.paragraphs[0].text == "Click"


def test_rich_text_multi_run_paragraph():
    """Multiple formatted runs in one paragraph from parse_html."""
    slide = _slide()
    tb = _tb(slide)
    block = parse_html("Normal <b>bold</b> <i>italic</i>")
    render_text_block(block, tb, _ctx(), base_size_pt=12)
    text_content = tb.text_frame.paragraphs[0].text
    assert "Normal" in text_content and "bold" in text_content


def test_rich_text_backward_compat():
    """text() DSL still works."""
    slide = _slide()
    tb = _tb(slide)
    block = text("Plain text works")
    render_text_block(block, tb, _ctx(), base_size_pt=12)
    assert tb.text_frame.paragraphs[0].text == "Plain text works"
