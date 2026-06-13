"""Edge cases for SlidePalette and palette application helpers."""

from __future__ import annotations

import pytest
from pptx import Presentation
from pptx.util import Inches

from slides_factory.palette import SlidePalette, apply_shape_palette_text
from slides_factory.palette import apply_paragraph_color


def test_usable_extras_honors_start_and_end():
    palette = SlidePalette(
        text="#000000",
        highlight="#111111",
        main=("#222222",),
        extras=("#333333", "#444444", "#555555", "#666666"),
        extras_start=1,
        extras_end=3,
    )
    assert palette.usable_extras() == ("#444444", "#555555")


def test_usable_extras_defaults_to_full_pool():
    palette = SlidePalette(
        text="#000000",
        highlight="#111111",
        main=("#222222",),
        extras=("#AAAAAA", "#BBBBBB"),
    )
    assert palette.usable_extras() == ("#AAAAAA", "#BBBBBB")


def test_extra_at_wraps_index():
    palette = SlidePalette(
        text="#000000",
        highlight="#111111",
        main=("#222222",),
        extras=("#AAAAAA", "#BBBBBB"),
    )
    assert palette.extra_at(0) == "#AAAAAA"
    assert palette.extra_at(2) == "#AAAAAA"


def test_extra_at_raises_when_pool_empty():
    palette = SlidePalette(
        text="#000000",
        highlight="#111111",
        main=("#222222",),
        extras=("#AAAAAA",),
        extras_start=1,
        extras_end=1,
    )
    with pytest.raises(ValueError, match="no usable extras"):
        palette.extra_at(0)


def test_apply_shape_palette_text_skips_shapes_without_text_frame():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_shape(1, Inches(1), Inches(1), Inches(1), Inches(1))
    palette = SlidePalette(
        text="#FF0000",
        highlight="#00FF00",
        main=("#0000FF",),
        extras=("#FFFFFF",),
    )
    apply_shape_palette_text(shape, palette)


def test_apply_paragraph_color_sets_runs():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    title.text = "Hello"
    apply_paragraph_color(title.text_frame.paragraphs[0], "#112233")
    assert str(title.text_frame.paragraphs[0].runs[0].font.color.rgb) == "112233"
