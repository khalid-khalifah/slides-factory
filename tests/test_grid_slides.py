"""Tests for GridSlideService and document grid-builder facade."""

from __future__ import annotations

from pathlib import Path

import pytest
from pptx.enum.shapes import MSO_SHAPE_TYPE

from slides_factory import document
from slides_factory.core.grid import GridSlideService
from slides_factory.core.manager import SlideManager
from slides_factory.core.session import PresentationSession


@pytest.fixture
def prs():
    theme_path = Path("themes/default.pptx")
    session = PresentationSession.create_new(
        default_theme_provider=lambda: theme_path,
        delete_slide_callback=lambda p, i: SlideManager(p).delete_slide(i),
    )
    return session.presentation


def _spec() -> dict:
    return {
        "grid": "grid-cols-2 gap-4",
        "cells": [
            {
                "element": {
                    "kind": "text",
                    "style": {"text_size": "lg", "bold": True},
                    "props": {"text": "Highlights"},
                }
            },
            {
                "element": {
                    "kind": "card",
                    "style": {"background_color": "surface"},
                    "props": {"title": "Revenue", "value": "$1.2M"},
                }
            },
        ],
    }


# --- GridSlideService -------------------------------------------------------


def test_new_grid_slide_starts_empty(prs):
    svc = GridSlideService(prs)
    result = svc.new_grid_slide(grid="grid-cols-2 gap-4")
    assert result["kind"] == "grid"
    assert result["data"]["grid"] == "grid-cols-2 gap-4"
    assert result["data"]["cells"] == []


def test_add_set_remove_cell_round_trip(prs):
    svc = GridSlideService(prs)
    svc.new_grid_slide(grid="grid-cols-2 grid-rows-2 gap-4")

    add = svc.add_cell(0, kind="text", at="col-span-2", props={"text": "Hello"})
    assert add["cell_index"] == 0

    svc.add_cell(0, kind="card", props={"title": "Rev", "value": "$1M"})
    spec = svc.require_grid_data(0)
    assert [c["element"]["kind"] for c in spec["cells"]] == ["text", "card"]
    assert spec["cells"][0]["at"] == "col-span-2"

    svc.set_cell(0, 0, props={"text": "Updated"})
    spec = svc.require_grid_data(0)
    assert spec["cells"][0]["element"]["props"]["text"] == "Updated"

    svc.remove_cell(0, 0)
    spec = svc.require_grid_data(0)
    assert [c["element"]["kind"] for c in spec["cells"]] == ["card"]


def test_set_slide_merges_frame_info_and_updates_grid(prs):
    svc = GridSlideService(prs)
    svc.new_grid_slide(grid="grid-cols-1", frame_info={"title": "Original"})

    svc.set_slide(0, grid="grid-cols-[2_1]", frame_info={"title": "Q3"})
    spec = svc.require_grid_data(0)
    assert spec["grid"] == "grid-cols-[2_1]"
    assert spec["frame_info"]["title"] == "Q3"


def test_add_cell_rejects_unknown_kind(prs):
    svc = GridSlideService(prs)
    svc.new_grid_slide()
    with pytest.raises(KeyError):
        svc.add_cell(0, kind="nope", props={})


def test_add_cell_requires_grid_slide(prs):
    svc = GridSlideService(prs)
    document.add_slide(
        prs, "simple", {"headline": {"text": "Plain"}, "body": {"text": "x"}}
    )
    with pytest.raises(ValueError, match="not a raw grid slide"):
        svc.add_cell(0, kind="text", props={"text": "x"})


def test_layout_blocked_on_cover_frame(prs, minimal_brand_yaml, tmp_path):
    output = tmp_path / "blocked.pptx"
    document.create_document(output, brand=minimal_brand_yaml)
    prs = document.open_document(output)
    svc = GridSlideService(prs)
    with pytest.raises(ValueError, match="does not allow grid layout"):
        svc.add_layout_slide({"grid": "grid-cols-1", "cells": []}, frame="cover")


def test_layout_slide_round_trip_without_brand(tmp_path):
    output = tmp_path / "grid.pptx"
    prs = document.create_document(output)
    result = document.add_layout_slide(prs, _spec())
    document.save_document(prs, output)

    assert result["kind"] == "grid"

    prs = document.open_document(output)
    info = document.get_slide_info(prs, 0)
    assert info["kind"] == "grid"
    assert info["data"]["grid"] == "grid-cols-2 gap-4"
    assert info["data"]["cells"][0]["element"]["props"]["text"] == "Highlights"

    slide = prs.slides[0]
    has_textbox = any(s.has_text_frame and not s.is_placeholder for s in slide.shapes)
    has_card = any(s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE for s in slide.shapes)
    assert has_textbox and has_card


def test_layout_slide_with_frame_draws_info_layer(tmp_path, minimal_brand_yaml):
    output = tmp_path / "branded.pptx"
    prs = document.create_document(output, brand=minimal_brand_yaml)
    spec = _spec()
    spec["frame_info"] = {"title": "Quarterly Review", "page_number": 1}
    result = document.add_layout_slide(prs, spec, frame="paneled")
    document.save_document(prs, output)

    assert result["frame_id"] == "paneled"

    prs = document.open_document(output)
    slide = prs.slides[0]
    texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
    assert "Quarterly Review" in texts


def test_layout_cells_land_inside_frame_playground(tmp_path, minimal_brand_yaml):
    output = tmp_path / "branded.pptx"
    prs = document.create_document(output, brand=minimal_brand_yaml)
    spec = {
        "frame_info": {"title": "T"},
        "grid": "grid-cols-1",
        "cells": [
            {
                "element": {
                    "kind": "card",
                    "style": {"background_color": "surface"},
                    "props": {"value": "X"},
                }
            }
        ],
    }
    document.add_layout_slide(prs, spec, frame="paneled")

    slide = prs.slides[0]
    card = next(s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE)
    width = int(prs.slide_width)
    height = int(prs.slide_height)
    assert abs(card.left - int(width * 0.10)) < width * 0.02
    assert abs(card.top - int(height * 0.25)) < height * 0.02


def test_document_facade_delegates_to_grid_service(prs):
    result = document.new_grid_slide(prs, grid="grid-cols-1")
    assert result["kind"] == "grid"
    document.add_cell(prs, 0, kind="text", props={"text": "Hi"})
    info = document.get_slide_info(prs, 0)
    assert info["data"]["cells"][0]["element"]["props"]["text"] == "Hi"
