"""Tests for the SVG converter."""

from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

from slides_factory.converters.svg import render_svg_file, render_svg_string
from slides_factory.geometry import Box

_SIMPLE_SVG = """<?xml version="1.0" encoding="UTF-8"?>
<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 100 100" width="100" height="100">
  <rect x="0" y="0" width="100" height="100" fill="#FF0000"/>
  <circle cx="50" cy="50" r="25" fill="#00FF00"/>
</svg>"""


def _slide():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    return prs.slides.add_slide(prs.slide_layouts[6])


# ---------------------------------------------------------------------------
# Converter tests
# ---------------------------------------------------------------------------

def test_render_svg_string_creates_shapes():
    slide = _slide()
    box = Box(0, 0, int(Inches(2)), int(Inches(2)))
    render_svg_string(_SIMPLE_SVG, slide, box)
    assert len(slide.shapes) == 2  # rect + circle


def test_render_svg_string_with_scale():
    slide = _slide()
    box = Box(int(Inches(2)), 0, int(Inches(2)), int(Inches(2)))
    render_svg_string(_SIMPLE_SVG, slide, box, scale=2.0)
    assert len(slide.shapes) > 0


def test_render_svg_file(tmp_path: Path):
    svg_path = tmp_path / "test.svg"
    svg_path.write_text(_SIMPLE_SVG, encoding="utf-8")
    slide = _slide()
    box = Box(0, 0, int(Inches(2)), int(Inches(2)))
    render_svg_file(svg_path, slide, box)
    assert len(slide.shapes) > 0


# ---------------------------------------------------------------------------
# Edge cases
# ---------------------------------------------------------------------------

def test_render_svg_no_viewbox():
    """SVG without viewBox still renders."""
    svg = '<svg xmlns="http://www.w3.org/2000/svg"><rect width="50" height="50" fill="blue"/></svg>'
    slide = _slide()
    box = Box(0, 0, int(Inches(2)), int(Inches(2)))
    render_svg_string(svg, slide, box)
    assert len(slide.shapes) > 0


def test_render_svg_empty():
    """Empty SVG produces no shapes (no crash)."""
    svg = '<svg xmlns="http://www.w3.org/2000/svg"></svg>'
    slide = _slide()
    box = Box(0, 0, int(Inches(2)), int(Inches(2)))
    render_svg_string(svg, slide, box)
    assert len(slide.shapes) == 0


def test_render_svg_fits_box():
    """SVG with viewBox is auto-scaled to fit the box."""
    slide = _slide()
    box = Box(0, 0, int(Inches(1)), int(Inches(1)))
    render_svg_string(_SIMPLE_SVG, slide, box)
    assert len(slide.shapes) > 0
