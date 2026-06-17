"""Tests for single-slide preview rendering."""

from __future__ import annotations

from pptx import Presentation

from slides_factory.preview import render as preview


def test_render_preview_pptx_produces_valid_deck(minimal_brand_yaml, tmp_path):
    data = {"headline": {"text": "Preview"}, "body": {"text": "Hello"}}
    pptx_bytes = preview.render_preview_pptx(
        "simple",
        data,
        brand=minimal_brand_yaml,
        frame="plain",
        locale="en",
    )
    assert pptx_bytes[:2] == b"PK"

    deck_path = tmp_path / "preview.pptx"
    deck_path.write_bytes(pptx_bytes)
    prs = Presentation(str(deck_path))
    assert len(prs.slides) == 1


def test_new_presentation_is_in_memory(minimal_brand_yaml):
    from slides_factory import document

    prs = document.new_presentation(brand=minimal_brand_yaml, locale="en")
    assert len(prs.slides) == 0
